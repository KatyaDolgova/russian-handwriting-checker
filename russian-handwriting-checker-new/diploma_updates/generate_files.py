"""
Генерирует обновлённые файлы для дипломного проекта:
  - updated_slides.pptx   — только изменённые слайды (6, 11, 12, 15, 22)
  - updated_report.docx   — только изменённые разделы отчёта
  - use_case_diagram.xml  — диаграмма вариантов использования (draw.io)
  - class_diagram.xml     — диаграмма классов (draw.io)
"""

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor, Inches as DInches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import os

OUT = os.path.dirname(__file__)

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
SLATE800 = RGBColor(0x1E, 0x29, 0x3B)
SLATE500 = RGBColor(0x64, 0x74, 0x8B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
EMERALD  = RGBColor(0x05, 0x96, 0x69)
RED      = RGBColor(0xDC, 0x26, 0x26)
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, l, t, w, h, fill=None, line=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                size=12, bold=False, color=None, align=PP_ALIGN.LEFT,
                wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or SLATE800
    return txBox


def add_card(slide, l, t, w, h, title, lines, title_color=INDIGO, bg=LIGHT_BG):
    add_rect(slide, l, t, w, h, fill=bg, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_textbox(slide, title, l+0.12, t+0.12, w-0.24, 0.35,
                size=11, bold=True, color=title_color)
    body = '\n'.join(lines)
    add_textbox(slide, body, l+0.12, t+0.48, w-0.24, h-0.58,
                size=9, color=SLATE500)


def slide_title(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.0, fill=INDIGO)
    add_textbox(slide, title, 0.3, 0.15, 12.0, 0.55,
                size=20, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle, 0.3, 0.65, 12.0, 0.3,
                    size=11, color=RGBColor(0xC7, 0xD2, 0xFE))


# ─────────────────────────────────────────────────────────────────────────────
# PPTX — обновлённые слайды
# ─────────────────────────────────────────────────────────────────────────────
def make_pptx():
    prs = new_prs()

    # ── СЛАЙД 6: Типы проверки ──────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFC))
    slide_title(sl, 'Типы проверки: встроенные функции',
                '6 функций готовы сразу после запуска системы')

    cards = [
        ('Орфография\nи пунктуация',
         ['Ошибки с объяснением', 'Оценка 0–5', 'Общий комментарий']),
        ('Сочинение\nОГЭ 9.3',
         ['8 критериев:', 'СК1–СК4 (содержание)', 'ГК1–ГК4 (грамотность)']),
        ('Сочинение\nЕГЭ 2026 ★',
         ['10 критериев К1–К10', 'по нормам ФИПИ 2026', 'Максимум 22 балла']),
        ('Итоговое сочинение\n(школьная)',
         ['5 критериев К1–К5', 'Зачёт / незачёт', 'Допуск к ЕГЭ']),
        ('Итоговое сочинение\n(вузовская)',
         ['10 критериев К1–К10', 'Максимум 20 баллов', 'Для поступления']),
        ('Свободная\nпроверка',
         ['Учитель сам', 'формулирует задание', 'ИИ его выполняет']),
    ]
    col_w, col_h = 2.05, 2.2
    col_gap = 0.12
    start_x, start_y = 0.22, 1.15
    for i, (title, lines) in enumerate(cards):
        x = start_x + i * (col_w + col_gap)
        add_card(sl, x, start_y, col_w, col_h, title, lines)

    add_rect(sl, 0.22, 3.6, 12.89, 0.5,
             fill=RGBColor(0xEE, 0xF2, 0xFF),
             line=RGBColor(0xC7, 0xD2, 0xFE))
    add_textbox(sl,
        '✚  Пользователь может создать собственную функцию, '
        'опубликовать в галерею и поделиться с другими учителями',
        0.4, 3.65, 12.5, 0.38, size=10, color=INDIGO)

    # ── СЛАЙД 11: CheckService ───────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFC))
    slide_title(sl, 'Как работает проверка: CheckService',
                'Routes → CheckService → LLMService → parse JSON → highlight errors')

    # Схема стрелок
    steps = [
        ('POST\n/check/stream', 0.2),
        ('FunctionRepository\n.get(id)', 2.1),
        ('_build_messages()\nsystem + user', 4.0),
        ('LLMService\n.stream()', 5.9),
        ('_parse()\nJSON → dict\n+очистка ```json```', 7.8),
        ('_build_result()', 9.9),
    ]
    for label, x in steps:
        add_card(sl, x, 1.1, 1.75, 1.0, label, [], title_color=SLATE800)
        if x < 9.9:
            add_textbox(sl, '→', x+1.75, 1.4, 0.35, 0.4, size=14, color=SLATE500)

    # Таблица результата
    add_textbox(sl, '_build_result() — что входит в итоговый ответ:',
                0.2, 2.35, 7.0, 0.35, size=11, bold=True, color=SLATE800)

    rows = [
        ('corrected_text', 'Исправленный текст'),
        ('errors',         'Список ошибок (JSON-массив)'),
        ('score / score_max', 'Числовая оценка и максимум'),
        ('pass_fail',      '«зачёт» / «незачёт» (итоговое сочинение)  ★ НОВОЕ'),
        ('criteria',       'Баллы по критериям ОГЭ/ЕГЭ  ★ НОВОЕ'),
        ('min_words',      'Минимум слов из настроек функции  ★ НОВОЕ'),
        ('html_highlighted', 'Текст с <span> вокруг ошибок'),
        ('comment',        'Общий комментарий'),
    ]
    row_h = 0.33
    for i, (field, desc) in enumerate(rows):
        y = 2.75 + i * row_h
        bg = RGBColor(0xEE,0xF2,0xFF) if 'НОВОЕ' in desc else WHITE
        add_rect(sl, 0.2, y, 2.8, row_h-0.02, fill=bg,
                 line=RGBColor(0xE2,0xE8,0xF0))
        add_rect(sl, 3.0, y, 6.5, row_h-0.02, fill=WHITE,
                 line=RGBColor(0xE2,0xE8,0xF0))
        add_textbox(sl, field, 0.3, y+0.04, 2.6, row_h-0.06,
                    size=8, bold=True, color=INDIGO)
        add_textbox(sl, desc, 3.1, y+0.04, 6.3, row_h-0.06,
                    size=8, color=SLATE800)

    # Код _parse
    code = (
        "def _parse(self, raw: str) -> dict:\n"
        "    # Очищаем markdown-обёртку ```json ... ```\n"
        "    cleaned = re.sub(r'^```(?:json)?\\s*', '', raw.strip())\n"
        "    cleaned = re.sub(r'\\s*```\\s*$', '', cleaned)\n"
        "    match = re.search(r'\\{[\\s\\S]*\\}', cleaned)\n"
        "    ..."
    )
    add_rect(sl, 9.7, 2.35, 3.45, 2.0,
             fill=RGBColor(0x1E,0x29,0x3B),
             line=RGBColor(0x33,0x41,0x55))
    add_textbox(sl, code, 9.8, 2.42, 3.25, 1.85,
                size=7, color=RGBColor(0x93,0xC5,0xFD))

    # ── СЛАЙД 12: Версионирование ────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFC))
    slide_title(sl, 'Функции проверки: создание, публикация, галерея',
                'Каждая функция — это system_prompt + user_template с плейсхолдером {text}')

    blocks = [
        ('Создание', INDIGO,
         ['Учитель пишет system_prompt', '(инструкции для ИИ)', 'и user_template с {text}']),
        ('Использование', INDIGO,
         ['Выбирается в CheckPanel', 'Текст подставляется вместо {text}', 'Отправляется в LLM']),
        ('Публикация', EMERALD,
         ['Флаг is_published = True', 'Функция попадает в галерею', 'Другие видят её версии']),
        ('Версионирование ★', AMBER,
         ['Авто: при редактировании опубл.', 'Вручную: кнопка «Сохранить версию»',
          'Удаление версий', 'Другие видят версии (чтение)']),
        ('Копирование', INDIGO,
         ['Другой учитель нажимает', '«Скопировать»', 'Копия в его библиотеке']),
    ]
    bw = 2.4
    gap = 0.17
    sx = 0.22
    for i, (title, color, lines) in enumerate(blocks):
        x = sx + i*(bw+gap)
        add_card(sl, x, 1.15, bw, 2.6, title, lines, title_color=color)
        if i < len(blocks)-1:
            add_textbox(sl, '→', x+bw+0.01, 2.1, 0.18, 0.5,
                        size=16, color=SLATE500, align=PP_ALIGN.CENTER)

    # Новые API-маршруты для версий
    add_rect(sl, 0.22, 4.0, 12.89, 1.3,
             fill=RGBColor(0xFF,0xFB,0xEB),
             line=RGBColor(0xFD,0xE6,0x8A))
    add_textbox(sl, '★  Новые маршруты версионирования:',
                0.4, 4.05, 12.5, 0.32, size=10, bold=True, color=AMBER)
    routes = (
        'POST   /functions/{id}/versions              — сохранить версию вручную (+ необязательная заметка)\n'
        'DELETE /functions/{id}/versions/{ver_id}     — удалить версию\n'
        'GET    /functions/{id}/versions              — история версий (владелец + все пользователи для опубликованных)'
    )
    add_textbox(sl, routes, 0.4, 4.38, 12.5, 0.82, size=9, color=SLATE800)

    add_rect(sl, 0.22, 5.45, 12.89, 0.42,
             fill=RGBColor(0xFF,0xF1,0xF2),
             line=RGBColor(0xFE,0xCA,0xCA))
    add_textbox(sl,
        '⚠  is_default = True → нельзя редактировать и удалять  '
        '(6 стандартных функций защищены от изменений)',
        0.4, 5.5, 12.5, 0.32, size=10, color=RED)

    # ── СЛАЙД 15: База данных ────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFC))
    slide_title(sl, 'База данных: 8 таблиц, PostgreSQL через Supabase', '')

    tables = [
        ('user_profiles',
         'user_id (PK), display_name, bio',
         'Профиль учителя'),
        ('checks  ★',
         'id, user_id, pupil_id, function_id,\nfolder_id, title, source_text,\noriginal_text, corrected_text,\nerrors (JSON), criteria (JSONB) ★,\npass_fail (TEXT) ★, score, score_max,\ncomment, work_date',
         'Результаты проверок'),
        ('functions  ★',
         'id, user_id, name, system_prompt,\nuser_template,\nscore_max (INT) ★,\nmin_words (INT) ★,\nis_default, is_published',
         'Функции проверки'),
        ('function_versions',
         'id, function_id, version_number,\nname, system_prompt, user_template,\nchange_note, created_at',
         'История версий'),
        ('pupils',
         'id, user_id, name\nUNIQUE(user_id, name)',
         'Ученики'),
        ('groups',
         'id, user_id, name, description',
         'Группы / классы'),
        ('pupil_groups',
         'pupil_id (FK), group_id (FK)',
         'Ученик ↔ Группа M:M'),
        ('folders',
         'id, user_id, name, description',
         'Папки для работ'),
    ]

    positions = [
        (0.15, 1.1), (2.55, 1.1), (5.25, 1.1), (8.25, 1.1),
        (0.15, 4.05), (2.55, 4.05), (5.25, 4.05), (8.25, 4.05),
    ]
    widths = [2.3, 2.6, 2.9, 4.9,  2.3, 2.6, 2.9, 4.9]
    heights = [1.3, 2.85, 2.85, 2.85,  1.3, 1.3, 1.3, 1.3]

    for i, (name, cols, desc) in enumerate(tables):
        x, y = positions[i]
        w = widths[i]
        h = heights[i]
        is_updated = '★' in name
        bg = RGBColor(0xEE,0xF2,0xFF) if is_updated else WHITE
        add_card(sl, x, y, w, h, name, [cols, '', desc],
                 title_color=INDIGO if is_updated else SLATE800, bg=bg)

    note = ('★  checks.criteria — JSONB: баллы по критериям {«К1»: {score, max, comment}}\n'
            '   checks.pass_fail — TEXT: «зачёт»/«незачёт» (исключается из числовой статистики)\n'
            '   functions.score_max, min_words — параметры функции для оценивания и требований к объёму')
    add_rect(sl, 0.15, 5.5, 13.0, 0.85,
             fill=RGBColor(0xEE,0xF2,0xFF),
             line=RGBColor(0xC7,0xD2,0xFE))
    add_textbox(sl, note, 0.28, 5.55, 12.7, 0.75, size=8, color=INDIGO)

    # ── СЛАЙД 22: Результаты ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8, 0xFA, 0xFC))
    slide_title(sl, 'Результаты работы', '')

    results = [
        ('🖊', 'Гибридный OCR',
         'PaddleOCR + Tesseract, автовыбор лучшего;\nподдержка кириллических имён файлов'),
        ('📄', 'Документы',
         'MarkItDown: PDF, DOCX, PPTX, TXT\nс сохранением структуры'),
        ('🤖', '6 типов проверки  ★',
         'Орфография, ОГЭ 9.3, ЕГЭ 2026,\nИт.сочинение (школьн./вузовск.), Свободная'),
        ('⚡', 'Стриминг SSE',
         'Ответ ИИ появляется в реальном времени,\nне надо ждать'),
        ('🔐', 'JWT-авторизация',
         'Supabase Auth,\nизоляция данных по user_id'),
        ('👨‍🎓', 'Управление данными',
         'Ученики, группы, папки, история;\nзачёт/незачёт в статистике  ★'),
        ('⚙️', 'Галерея функций',
         'Публикация, копирование, версионирование,\nудаление версий, просмотр другими  ★'),
        ('🔍', 'Подсветка ошибок в истории  ★',
         'Волнистая линия + тултип:\nтип, исправление, пояснение'),
        ('🐳', 'Docker',
         'Готово к деплою\nна любом сервере'),
    ]

    cw, ch = 4.0, 1.25
    gap = 0.13
    sx = 0.22
    cols = 3
    for i, (icon, title, desc) in enumerate(results):
        row = i // cols
        col = i % cols
        x = sx + col*(cw+gap)
        y = 1.1 + row*(ch+gap)
        is_new = '★' in title or '★' in desc
        bg = RGBColor(0xEE,0xF2,0xFF) if is_new else WHITE
        tc = INDIGO if is_new else SLATE800
        add_rect(sl, x, y, cw, ch, fill=bg, line=RGBColor(0xE2,0xE8,0xF0))
        add_textbox(sl, icon+' '+title.replace('  ★',''), x+0.1, y+0.08,
                    cw-0.2, 0.38, size=11, bold=True, color=tc)
        add_textbox(sl, desc, x+0.1, y+0.46, cw-0.2, ch-0.55,
                    size=9, color=SLATE500)

    add_textbox(sl, 'Спасибо за внимание!',
                0, 6.9, 13.33, 0.5, size=14, bold=True,
                color=INDIGO, align=PP_ALIGN.CENTER)

    out = os.path.join(OUT, 'updated_slides.pptx')
    prs.save(out)
    print(f'✓ {out}')


# ─────────────────────────────────────────────────────────────────────────────
# DOCX — обновлённые разделы отчёта
# ─────────────────────────────────────────────────────────────────────────────
def set_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in h.runs:
        run.font.color.rgb = DRGBColor(0x4F, 0x46, 0xE5)
    return h


def add_label(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.bold = True
    run.font.color.rgb = DRGBColor(0x4F, 0x46, 0xE5)
    run.font.size = DPt(11)
    p.paragraph_format.space_before = DPt(8)
    return p


def add_note(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.italic = True
    run.font.color.rgb = DRGBColor(0x64, 0x74, 0x8B)
    run.font.size = DPt(9)
    p.paragraph_format.left_indent = Cm(0.5)
    return p


def add_body(doc, text):
    p = doc.add_paragraph(text)
    p.paragraph_format.first_line_indent = Cm(1.25)
    p.paragraph_format.space_after = DPt(4)
    for run in p.runs:
        run.font.size = DPt(12)
    return p


def add_code_block(doc, code):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Cm(1)
    p.paragraph_format.space_before = DPt(4)
    p.paragraph_format.space_after = DPt(4)
    run = p.add_run(code)
    run.font.name = 'Courier New'
    run.font.size = DPt(9)
    run.font.color.rgb = DRGBColor(0x1E, 0x29, 0x3B)
    # light gray bg through shading
    pPr = p._p.get_or_add_pPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), 'F1F5F9')
    pPr.append(shd)
    return p


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1+len(rows), cols=len(headers))
    table.style = 'Table Grid'
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = h
        for run in cell.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = DPt(10)
            run.font.color.rgb = DRGBColor(0xFF, 0xFF, 0xFF)
        cell._tc.get_or_add_tcPr().append(
            OxmlElement('w:shd'))
        shd = cell._tc.tcPr.findall(qn('w:shd'))[-1]
        shd.set(qn('w:val'), 'clear')
        shd.set(qn('w:color'), 'auto')
        shd.set(qn('w:fill'), '4F46E5')

    for ri, row_data in enumerate(rows):
        row = table.rows[ri+1]
        is_new = any('★' in str(c) for c in row_data)
        for ci, cell_text in enumerate(row_data):
            cell = row.cells[ci]
            cell.text = str(cell_text)
            for run in cell.paragraphs[0].runs:
                run.font.size = DPt(9)
                if is_new:
                    run.font.color.rgb = DRGBColor(0x4F, 0x46, 0xE5)
            if is_new:
                shd = OxmlElement('w:shd')
                shd.set(qn('w:val'), 'clear')
                shd.set(qn('w:color'), 'auto')
                shd.set(qn('w:fill'), 'EEF2FF')
                cell._tc.get_or_add_tcPr().append(shd)
    return table


def make_docx():
    doc = Document()

    # Стиль по умолчанию
    style = doc.styles['Normal']
    style.font.name = 'Times New Roman'
    style.font.size = DPt(12)

    # Заголовок документа
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run('ОБНОВЛЁННЫЕ РАЗДЕЛЫ ОТЧЁТА\nРусЯзык AI — дипломный проект')
    run.font.bold = True
    run.font.size = DPt(14)
    run.font.color.rgb = DRGBColor(0x4F, 0x46, 0xE5)

    doc.add_paragraph()
    add_note(doc, '★ — новые или изменённые элементы. Вставь этот текст вместо соответствующих абзацев в оригинальном отчёте.')
    doc.add_paragraph()

    # ── 1.3 Требования ──────────────────────────────────────────────────────
    set_heading(doc, '1.3 Требования к системе', level=2)
    add_label(doc, 'Пункт 5 функциональных требований — ЗАМЕНИТЬ:')
    add_body(doc,
        '5) результат проверки должен включать исправленный вариант текста, список ошибок '
        'с подсветкой и пояснениями, числовую оценку, а для работ, проверяемых по критериям '
        'ОГЭ/ЕГЭ — баллы по каждому критерию отдельно; для итогового сочинения вместо '
        'числовой оценки должен выводиться результат «зачёт» / «незачёт»;')

    add_label(doc, 'Новый пункт 10 — ДОБАВИТЬ после пункта 9:')
    add_body(doc,
        '10) система должна поддерживать ручное сохранение версий функции проверки с '
        'необязательной заметкой об изменении; версии должны быть доступны для просмотра '
        'другими пользователями для опубликованных функций; должна быть предусмотрена '
        'возможность удаления отдельных версий.')

    add_label(doc, 'Пункт нефункциональных требований — ДОБАВИТЬ:')
    add_body(doc,
        '7) система должна корректно обрабатывать файлы с кириллическими именами на '
        'платформе Windows.')

    # ── 3.3 Диаграмма классов ───────────────────────────────────────────────
    set_heading(doc, '3.3 Диаграмма классов', level=2)
    add_label(doc, 'Абзац про модель Check — ЗАМЕНИТЬ:')
    add_body(doc,
        'Модель Check представляет результат одной проверки и содержит следующие атрибуты: '
        'идентификатор записи, идентификаторы пользователя, ученика, функции и папки, '
        'имя загруженного файла, название работы (title), исходный текст задания '
        '(source_text — для диктантов и изложений, где требуется сравнение с эталоном), '
        'текст ученика (original_text), исправленный текст (corrected_text), массив ошибок '
        '(errors, тип JSON), баллы по критериям (criteria, тип JSONB — объект вида '
        '{«К1»: {score, max, comment}}), результат зачёт/незачёт (pass_fail, тип TEXT), '
        'числовая оценка (score), максимальный балл (score_max), текстовый комментарий, '
        'дата работы (work_date) и дата создания записи (created_at).')
    add_body(doc,
        'Поля criteria и pass_fail являются взаимоисключающими: при наличии pass_fail запись '
        'исключается из расчёта числовой статистики и отображается отдельным счётчиком '
        'зачётов в интерфейсе.')

    add_label(doc, 'Абзац про модель Function — ЗАМЕНИТЬ:')
    add_body(doc,
        'Модель Function хранит настраиваемые функции проверки. Помимо идентификатора, '
        'имени, описания, системного промпта и шаблона пользователя, модель содержит два '
        'дополнительных числовых параметра: score_max (максимальный балл) и min_words '
        '(минимальное количество слов, необходимое для работы данного типа). Флаги '
        'is_default и is_published управляют доступностью функции.')

    add_label(doc, 'В описание FunctionRepository — ДОПОЛНИТЬ в конец:')
    add_body(doc,
        'Метод create_version(function_id, change_note) создаёт снапшот текущего состояния '
        'функции как новую запись в таблице function_versions с автоматически назначаемым '
        'номером версии. Это позволяет учителю вручную зафиксировать значимое изменение, '
        'не публикуя функцию в галерею. Метод delete_version(version_id) удаляет '
        'конкретную версию; текущая функция при этом не изменяется.')

    # ── 3.4 Схема базы данных ───────────────────────────────────────────────
    set_heading(doc, '3.4 Схема базы данных', level=2)
    add_label(doc, 'Таблица checks — обновлённый список столбцов:')

    add_table(doc,
        ['Столбец', 'Тип', 'Описание'],
        [
            ('id',             'VARCHAR (PK)',   'UUID записи'),
            ('user_id',        'VARCHAR',        'Идентификатор учителя (Supabase Auth)'),
            ('pupil_id',       'VARCHAR (FK)',   'Ученик (nullable)'),
            ('function_id',    'VARCHAR',        'Использованная функция'),
            ('folder_id',      'VARCHAR (FK)',   'Папка (nullable)'),
            ('filename',       'VARCHAR',        'Имя загруженного файла'),
            ('title',          'VARCHAR',        'Название работы (nullable)'),
            ('source_text',    'TEXT',           'Исходный текст задания (для диктантов)'),
            ('original_text',  'TEXT',           'Текст ученика'),
            ('corrected_text', 'TEXT',           'Исправленный вариант'),
            ('errors',         'JSON',           'Массив ошибок [{original, corrected, type, comment}]'),
            ('criteria ★',     'JSONB',          'Баллы по критериям {К1: {score, max, comment}}'),
            ('pass_fail ★',    'TEXT',           '«зачёт» / «незачёт» (nullable)'),
            ('score',          'FLOAT',          'Числовая оценка (nullable)'),
            ('score_max',      'FLOAT',          'Максимальный балл (nullable)'),
            ('comment',        'TEXT',           'Общий комментарий'),
            ('work_date',      'TIMESTAMP(TZ)',  'Дата выполнения работы (nullable)'),
            ('created_at',     'TIMESTAMP(TZ)',  'Дата создания записи (server default)'),
        ])

    doc.add_paragraph()
    add_label(doc, 'Таблица functions — обновлённый список столбцов:')
    add_table(doc,
        ['Столбец', 'Тип', 'Описание'],
        [
            ('id',                   'VARCHAR (PK)', 'UUID функции'),
            ('user_id',              'VARCHAR',      'Владелец (nullable для is_default)'),
            ('name',                 'VARCHAR',      'Название'),
            ('description',          'TEXT',         'Описание'),
            ('system_prompt',        'TEXT',         'Системный промпт для LLM'),
            ('user_template',        'TEXT',         'Шаблон с {text}'),
            ('score_max ★',          'INTEGER',      'Максимальный балл функции (nullable)'),
            ('min_words ★',          'INTEGER',      'Минимальное кол-во слов (nullable)'),
            ('is_default',           'BOOLEAN',      'Стандартная функция'),
            ('is_published',         'BOOLEAN',      'Опубликована в галерею'),
            ('original_function_id', 'VARCHAR',      'Источник при копировании (nullable)'),
        ])

    # ── 4.2 CheckService ────────────────────────────────────────────────────
    doc.add_paragraph()
    set_heading(doc, '4.2 Серверная часть: ключевые компоненты', level=2)
    add_label(doc, 'Добавить абзац про _parse() — после описания метода:')
    add_body(doc,
        'Метод _parse() выполняет разбор ответа языковой модели. Перед поиском JSON-блока '
        'метод очищает возможную Markdown-обёртку (```json ... ```), которую модель иногда '
        'добавляет к ответу вопреки инструкции. Очистка выполняется регулярными выражениями:')
    add_code_block(doc,
        "cleaned = re.sub(r'^```(?:json)?\\s*', '', raw.strip())\n"
        "cleaned = re.sub(r'\\s*```\\s*$', '', cleaned)\n"
        "match = re.search(r'\\{[\\s\\S]*\\}', cleaned)")
    add_body(doc,
        'После очистки метод ищет JSON-блок и при ошибке десериализации пытается исправить '
        'типичную ошибку — лишние запятые перед закрывающей скобкой. Параметр max_tokens '
        'увеличен до 8000 для корректной обработки больших текстов с множеством ошибок.')

    add_label(doc, 'Добавить абзац про _build_result():')
    add_body(doc,
        'Метод _build_result() формирует итоговый словарь результата. Помимо ранее описанных '
        'полей, метод теперь передаёт: criteria — словарь баллов по критериям из ответа LLM; '
        'pass_fail — результат «зачёт»/«незачёт»; score_max и min_words — из настроек функции; '
        'score_label — строковое представление оценки, если она не является числом. '
        'Флаг is_generation устанавливается в True, когда ответ не содержит полей «corrected» '
        'и «errors», то есть модель сгенерировала произвольный текст вместо структурированного JSON.')

    add_label(doc, 'Добавить подраздел 4.2.X — Поддержка кириллических имён файлов:')
    add_body(doc,
        'Функция cv2.imread() библиотеки OpenCV не поддерживает пути с символами кириллицы '
        'на платформе Windows из-за ограничений Windows API. Для решения этой проблемы '
        'реализована вспомогательная функция _imread_unicode(), которая читает файл как '
        'байтовый поток через numpy.fromfile() и передаёт буфер в cv2.imdecode():')
    add_code_block(doc,
        "def _imread_unicode(image_path: str) -> np.ndarray:\n"
        "    buf = np.fromfile(image_path, dtype=np.uint8)\n"
        "    img = cv2.imdecode(buf, cv2.IMREAD_COLOR)\n"
        "    if img is None:\n"
        "        raise ValueError(f'Не удалось загрузить: {image_path}')\n"
        "    return img")

    add_label(doc, 'Добавить абзац про версионирование функций — ЗАМЕНИТЬ:')
    add_body(doc,
        'Система поддерживает два режима версионирования. Автоматический режим: при сохранении '
        'изменений в уже опубликованной функции маршрут PUT /functions/{id} сохраняет значение '
        'флага is_published до вызова repo.update() и после обновления вызывает repo.republish(), '
        'который создаёт новую запись в function_versions. Ручной режим: учитель может в любой '
        'момент зафиксировать снапшот функции через POST /functions/{id}/versions с необязательным '
        'параметром change_note — это позволяет сохранять версии черновиков, не публикуя их. '
        'Версии можно удалять через DELETE /functions/{id}/versions/{ver_id}. История версий '
        'опубликованных функций доступна для просмотра другим пользователям (только чтение).')

    # ── 4.3 Клиентская часть ────────────────────────────────────────────────
    set_heading(doc, '4.3 Клиентская часть: ключевые компоненты', level=2)
    add_label(doc, 'Добавить описание компонента HighlightedText:')
    add_body(doc,
        'Компонент HighlightedText реализует интерактивную подсветку ошибок в сохранённом '
        'тексте проверки. Компонент принимает исходный текст и массив ошибок, разбивает '
        'строку на сегменты по вхождениям поля original каждой ошибки и рендерит каждый '
        'ошибочный сегмент как отдельный span с волнистым красным подчёркиванием '
        '(Tailwind-класс decoration-wavy decoration-red-400). При наведении курсора '
        'появляется тултип с тремя полями: тип ошибки, исправленный вариант (corrected) '
        'и краткое пояснение (comment). Тултип реализован через React-состояние '
        '(onMouseEnter / onMouseLeave) без внешних библиотек и имеет фиксированную ширину '
        '256 px с переносом длинного текста. Компонент используется в HistoryPanel для '
        'отображения исходного текста работы.')

    add_label(doc, 'Добавить описание зачёт/незачёт в ResultPanel:')
    add_body(doc,
        'Для функций с оценкой «зачёт»/«незачёт» компонент ResultPanel определяет тип '
        'результата по полю pass_fail или по наличию поля result в критериях. При '
        'обнаружении такого результата вместо числовых полей ввода оценки отображается '
        'выпадающий список с вариантами «Зачёт» и «Незачёт». В панели истории и на '
        'странице учеников записи с pass_fail отображаются отдельным счётчиком '
        '«X/Y зачётов» и не включаются в расчёт среднего процента успеваемости.')

    # ── Заключение ───────────────────────────────────────────────────────────
    set_heading(doc, 'Заключение — дополнение к перечню результатов', level=2)
    add_label(doc, 'ДОБАВИТЬ в список результатов:')
    items = [
        'реализована поддержка оценки «зачёт»/«незачёт» для итогового сочинения с отдельным '
        'отображением в статистике ученика и истории проверок;',
        'функции проверки приведены в соответствие с требованиями ФИПИ 2026: сочинение ЕГЭ '
        '(К1–К10, 22 балла), итоговое сочинение в двух вариантах — школьная оценка '
        '(зачёт/незачёт) и вузовская оценка (20 баллов);',
        'в истории проверок реализована подсветка ошибок с интерактивными тултипами;',
        'версионирование функций расширено: добавлены ручное сохранение версий, удаление '
        'версий и просмотр истории версий другими пользователями;',
        'устранена проблема чтения изображений с кириллическими именами файлов на платформе Windows.',
    ]
    for item in items:
        p = doc.add_paragraph(style='List Bullet')
        run = p.add_run(item)
        run.font.size = DPt(12)

    out = os.path.join(OUT, 'updated_report.docx')
    doc.save(out)
    print(f'✓ {out}')


# ─────────────────────────────────────────────────────────────────────────────
# XML — диаграммы draw.io
# ─────────────────────────────────────────────────────────────────────────────
USE_CASE_XML = '''<?xml version="1.0" encoding="UTF-8"?>
<mxGraphModel dx="1422" dy="762" grid="1" gridSize="10" guides="1"
  tooltips="1" connect="1" arrows="1" fold="1" page="1"
  pageScale="1" pageWidth="1654" pageHeight="1169" math="0" shadow="0">
  <root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>

    <!-- Актор -->
    <mxCell id="actor" value="Преподаватель" style="shape=mxgraph.flowchart.actor;whiteSpace=wrap;html=1;fontSize=14;fontStyle=1;" vertex="1" parent="1">
      <mxGeometry x="40" y="480" width="80" height="120" as="geometry"/>
    </mxCell>

    <!-- Граница системы -->
    <mxCell id="sys" value="РусЯзык AI" style="points=[[0,0],[0.25,0],[0.5,0],[0.75,0],[1,0],[1,0.25],[1,0.5],[1,0.75],[1,1],[0.75,1],[0.5,1],[0.25,1],[0,1],[0,0.75],[0,0.5],[0,0.25]];shape=mxgraph.flowchart.start_2;fillColor=none;strokeColor=#4F46E5;fontSize=16;fontStyle=1;align=center;verticalAlign=top;strokeWidth=3;" vertex="1" parent="1">
      <mxGeometry x="200" y="20" width="1400" height="1120" as="geometry"/>
    </mxCell>

    <!-- GROUP 1: Аутентификация -->
    <mxCell id="g1" value="Аутентификация и профиль" style="swimlane;fillColor=#EEF2FF;strokeColor=#4F46E5;fontSize=12;fontStyle=1;startSize=30;" vertex="1" parent="1">
      <mxGeometry x="220" y="40" width="340" height="340" as="geometry"/>
    </mxCell>
    <mxCell id="uc01" value="УД-01 Зарегистрироваться" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g1"><mxGeometry x="20" y="40" width="200" height="44" as="geometry"/></mxCell>
    <mxCell id="uc02" value="УД-02 Войти в систему" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g1"><mxGeometry x="20" y="94" width="200" height="44" as="geometry"/></mxCell>
    <mxCell id="uc03" value="УД-03 Выйти из системы" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g1"><mxGeometry x="20" y="148" width="200" height="44" as="geometry"/></mxCell>
    <mxCell id="uc04" value="УД-04 Изменить пароль" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g1"><mxGeometry x="20" y="202" width="200" height="44" as="geometry"/></mxCell>
    <mxCell id="uc05" value="УД-05 Редактировать профиль" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g1"><mxGeometry x="20" y="256" width="200" height="44" as="geometry"/></mxCell>

    <!-- GROUP 2: Проверка текста -->
    <mxCell id="g2" value="Проверка текста" style="swimlane;fillColor=#F0FDF4;strokeColor=#059669;fontSize=12;fontStyle=1;startSize=30;" vertex="1" parent="1">
      <mxGeometry x="580" y="40" width="380" height="580" as="geometry"/>
    </mxCell>
    <mxCell id="uc06" value="УД-06 Загрузить файл" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="40" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc07" value="УД-07 Ввести текст вручную" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="94" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc08" value="УД-08 Выбрать функцию проверки" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="148" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc09" value="УД-09 Запустить проверку текста" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;fontStyle=1;" vertex="1" parent="g2"><mxGeometry x="20" y="202" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc10" value="УД-10 Ввести исходный текст" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g2"><mxGeometry x="260" y="202" width="100" height="44" as="geometry"/></mxCell>
    <mxCell id="uc11" value="УД-11 Сохранить результат" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g2"><mxGeometry x="260" y="256" width="100" height="44" as="geometry"/></mxCell>
    <mxCell id="uc12" value="УД-12 Скопировать текст" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="256" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc13" value="УД-13 Распечатать отчёт" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="310" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc14" value="УД-14 Просмотреть историю" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;fontStyle=1;" vertex="1" parent="g2"><mxGeometry x="20" y="364" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc15" value="УД-15 Фильтровать историю" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g2"><mxGeometry x="260" y="364" width="100" height="44" as="geometry"/></mxCell>
    <mxCell id="uc16" value="УД-16 Просмотреть ошибки с тултипами ★" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#4F46E5;fillColor=#EEF2FF;" vertex="1" parent="g2"><mxGeometry x="260" y="418" width="100" height="50" as="geometry"/></mxCell>
    <mxCell id="uc17" value="УД-17 Редактировать запись" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g2"><mxGeometry x="20" y="418" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc18" value="УД-18 Удалить запись" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="472" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc19" value="УД-19 Переместить в папку (массово)" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g2"><mxGeometry x="20" y="526" width="220" height="44" as="geometry"/></mxCell>

    <!-- GROUP 3: Функции проверки -->
    <mxCell id="g3" value="Функции проверки" style="swimlane;fillColor=#FFF7ED;strokeColor=#F59E0B;fontSize=12;fontStyle=1;startSize=30;" vertex="1" parent="1">
      <mxGeometry x="980" y="40" width="400" height="620" as="geometry"/>
    </mxCell>
    <mxCell id="uc20" value="УД-20 Создать функцию" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g3"><mxGeometry x="20" y="40" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc21" value="УД-21 Редактировать функцию" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g3"><mxGeometry x="20" y="94" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc22" value="УД-22 Удалить функцию" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g3"><mxGeometry x="20" y="148" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc23" value="УД-23 Опубликовать функцию" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;fontStyle=1;" vertex="1" parent="g3"><mxGeometry x="20" y="202" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc24" value="УД-24 Снять с публикации" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g3"><mxGeometry x="260" y="202" width="120" height="44" as="geometry"/></mxCell>
    <mxCell id="uc25" value="УД-25 Сохранить версию вручную ★" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#4F46E5;fillColor=#EEF2FF;" vertex="1" parent="g3"><mxGeometry x="20" y="256" width="220" height="50" as="geometry"/></mxCell>
    <mxCell id="uc26" value="УД-26 Просмотреть историю версий" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;fontStyle=1;" vertex="1" parent="g3"><mxGeometry x="20" y="316" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc27" value="УД-27 Удалить версию ★" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#4F46E5;fillColor=#EEF2FF;" vertex="1" parent="g3"><mxGeometry x="260" y="316" width="120" height="44" as="geometry"/></mxCell>
    <mxCell id="uc28" value="УД-28 Просмотреть галерею" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;fontStyle=1;" vertex="1" parent="g3"><mxGeometry x="20" y="370" width="220" height="44" as="geometry"/></mxCell>
    <mxCell id="uc29" value="УД-29 Просмотреть версии чужой функции ★" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#4F46E5;fillColor=#EEF2FF;" vertex="1" parent="g3"><mxGeometry x="260" y="370" width="120" height="52" as="geometry"/></mxCell>
    <mxCell id="uc30" value="УД-30 Скопировать из галереи" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;strokeColor=#F59E0B;" vertex="1" parent="g3"><mxGeometry x="260" y="432" width="120" height="44" as="geometry"/></mxCell>

    <!-- GROUP 4: Ученики и группы -->
    <mxCell id="g4" value="Ученики и группы" style="swimlane;fillColor=#F0FEFF;strokeColor=#0891B2;fontSize=12;fontStyle=1;startSize=30;" vertex="1" parent="1">
      <mxGeometry x="220" y="400" width="340" height="360" as="geometry"/>
    </mxCell>
    <mxCell id="uc31" value="УД-31 Добавить ученика" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="40" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc32" value="УД-32 Редактировать ученика" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="90" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc33" value="УД-33 Удалить ученика" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="140" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc34" value="УД-34 Создать группу" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="190" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc35" value="УД-35 Управлять составом группы" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="240" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc36" value="УД-36 Просмотреть статистику ученика" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g4"><mxGeometry x="20" y="290" width="200" height="46" as="geometry"/></mxCell>

    <!-- GROUP 5: Папки -->
    <mxCell id="g5" value="Папки" style="swimlane;fillColor=#FFF1F2;strokeColor=#DC2626;fontSize=12;fontStyle=1;startSize=30;" vertex="1" parent="1">
      <mxGeometry x="220" y="780" width="340" height="220" as="geometry"/>
    </mxCell>
    <mxCell id="uc38" value="УД-38 Создать папку" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g5"><mxGeometry x="20" y="40" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc39" value="УД-39 Переименовать папку" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g5"><mxGeometry x="20" y="90" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc40" value="УД-40 Удалить папку" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g5"><mxGeometry x="20" y="140" width="200" height="40" as="geometry"/></mxCell>
    <mxCell id="uc41" value="УД-41 Переместить проверку в папку" style="ellipse;whiteSpace=wrap;html=1;fontSize=10;" vertex="1" parent="g5"><mxGeometry x="20" y="170" width="200" height="40" as="geometry"/></mxCell>

    <!-- Связи актор → группы (ассоциации) -->
    <mxCell id="e1" style="edgeStyle=orthogonalEdgeStyle;" edge="1" source="actor" target="g1" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="e2" style="edgeStyle=orthogonalEdgeStyle;" edge="1" source="actor" target="g2" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="e3" style="edgeStyle=orthogonalEdgeStyle;" edge="1" source="actor" target="g3" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="e4" style="edgeStyle=orthogonalEdgeStyle;" edge="1" source="actor" target="g4" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="e5" style="edgeStyle=orthogonalEdgeStyle;" edge="1" source="actor" target="g5" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>

    <!-- extend связи (пунктир) -->
    <mxCell id="ex1" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;" edge="1" source="uc09" target="uc10" parent="g2"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex2" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;" edge="1" source="uc09" target="uc11" parent="g2"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex3" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;" edge="1" source="uc14" target="uc15" parent="g2"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex4" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;strokeColor=#4F46E5;" edge="1" source="uc14" target="uc16" parent="g2"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex5" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;" edge="1" source="uc23" target="uc24" parent="g3"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex6" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;strokeColor=#4F46E5;" edge="1" source="uc26" target="uc27" parent="g3"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex7" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;strokeColor=#4F46E5;" edge="1" source="uc28" target="uc29" parent="g3"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="ex8" value="&lt;&lt;extend&gt;&gt;" style="edgeStyle=orthogonalEdgeStyle;dashed=1;fontSize=9;" edge="1" source="uc28" target="uc30" parent="g3"><mxGeometry relative="1" as="geometry"/></mxCell>

    <!-- Легенда -->
    <mxCell id="leg" value="★ — новые варианты использования (добавлены в текущей версии)" style="text;html=1;strokeColor=none;fillColor=#EEF2FF;align=left;verticalAlign=middle;spacingLeft=4;spacingRight=4;overflow=hidden;rotatable=0;fontSize=12;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="1">
      <mxGeometry x="40" y="1080" width="700" height="40" as="geometry"/>
    </mxCell>
  </root>
</mxGraphModel>'''


CLASS_XML = '''<?xml version="1.0" encoding="UTF-8"?>
<mxGraphModel dx="1422" dy="762" grid="1" gridSize="10" guides="1"
  tooltips="1" connect="1" arrows="1" fold="1" page="1"
  pageScale="1" pageWidth="1654" pageHeight="1169" math="0" shadow="0">
  <root>
    <mxCell id="0"/><mxCell id="1" parent="0"/>

    <!-- ═══════ DOMAIN MODELS ═══════ -->
    <mxCell id="hdr1" value="Доменные модели (таблицы БД)" style="text;html=1;strokeColor=none;fillColor=#4F46E5;fontColor=#ffffff;align=center;fontSize=13;fontStyle=1;spacingLeft=4;" vertex="1" parent="1">
      <mxGeometry x="20" y="20" width="380" height="34" as="geometry"/>
    </mxCell>

    <!-- Check -->
    <mxCell id="check" value="&lt;b&gt;Check&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#EEF2FF;strokeColor=#4F46E5;" vertex="1" parent="1">
      <mxGeometry x="20" y="64" width="380" height="460" as="geometry"/>
    </mxCell>
    <mxCell id="ch1" value="+ id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="26" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch2" value="+ user_id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="44" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch3" value="+ pupil_id: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="62" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch4" value="+ function_id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="80" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch5" value="+ folder_id: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="98" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch6" value="+ filename: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="116" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch7" value="+ title: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="134" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch8" value="+ source_text: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="152" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch9" value="+ original_text: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="170" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch10" value="+ corrected_text: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="188" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch11" value="+ errors: JSON | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="206" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch12" value="★ + criteria: JSONB | None" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="check"><mxGeometry x="0" y="224" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch13" value="★ + pass_fail: str | None" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="check"><mxGeometry x="0" y="242" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch14" value="+ score: float | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="260" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch15" value="+ score_max: float | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="278" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch16" value="+ comment: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="296" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch17" value="+ work_date: datetime | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="314" width="380" height="18" as="geometry"/></mxCell>
    <mxCell id="ch18" value="+ created_at: datetime" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="check"><mxGeometry x="0" y="332" width="380" height="18" as="geometry"/></mxCell>

    <!-- Function -->
    <mxCell id="func" value="&lt;b&gt;Function&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#EEF2FF;strokeColor=#4F46E5;" vertex="1" parent="1">
      <mxGeometry x="420" y="64" width="300" height="300" as="geometry"/>
    </mxCell>
    <mxCell id="fn1" value="+ id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="26" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn2" value="+ user_id: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="44" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn3" value="+ name: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="62" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn4" value="+ description: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="80" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn5" value="+ system_prompt: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="98" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn6" value="+ user_template: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="116" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn7" value="★ + score_max: int | None" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="func"><mxGeometry x="0" y="134" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn8" value="★ + min_words: int | None" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="func"><mxGeometry x="0" y="152" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn9" value="+ is_default: bool" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="170" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn10" value="+ is_published: bool" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="188" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fn11" value="+ original_function_id: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="func"><mxGeometry x="0" y="206" width="300" height="18" as="geometry"/></mxCell>

    <!-- FunctionVersion -->
    <mxCell id="fver" value="&lt;b&gt;FunctionVersion&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#F0FDF4;strokeColor=#059669;" vertex="1" parent="1">
      <mxGeometry x="420" y="374" width="300" height="210" as="geometry"/>
    </mxCell>
    <mxCell id="fv1" value="+ id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="26" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv2" value="+ function_id: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="44" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv3" value="+ version_number: int" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="62" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv4" value="+ name: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="80" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv5" value="+ system_prompt: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="98" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv6" value="+ user_template: str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="116" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv7" value="+ change_note: str | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="134" width="300" height="18" as="geometry"/></mxCell>
    <mxCell id="fv8" value="+ created_at: datetime" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fver"><mxGeometry x="0" y="152" width="300" height="18" as="geometry"/></mxCell>

    <!-- ═══════ SERVICES ═══════ -->
    <mxCell id="hdr2" value="Сервисы (бизнес-логика)" style="text;html=1;strokeColor=none;fillColor=#059669;fontColor=#ffffff;align=center;fontSize=13;fontStyle=1;spacingLeft=4;" vertex="1" parent="1">
      <mxGeometry x="740" y="20" width="360" height="34" as="geometry"/>
    </mxCell>

    <!-- CheckService -->
    <mxCell id="csvc" value="&lt;b&gt;CheckService&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#F0FDF4;strokeColor=#059669;" vertex="1" parent="1">
      <mxGeometry x="740" y="64" width="360" height="280" as="geometry"/>
    </mxCell>
    <mxCell id="cs1" value="- llm: LLMService" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="26" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs2" value="- function_repo: FunctionRepository" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="44" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="csd" value="──────────────────────────────" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=4;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="62" width="360" height="14" as="geometry"/></mxCell>
    <mxCell id="cs3" value="+ run_check(text, function_id) → dict" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="76" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs4" value="+ stream_check(text, function_id) → AsyncIterator" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="94" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs5" value="- _build_messages(text, function_id)" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="112" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs6" value="- _build_result(text, data, raw, score_max, min_words)" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="130" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs7" value="★ - _parse(raw) → dict  [очистка ```json```]" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="csvc"><mxGeometry x="0" y="148" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="cs8" value="- _highlight(text, errors) → str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="csvc"><mxGeometry x="0" y="166" width="360" height="18" as="geometry"/></mxCell>

    <!-- LLMService -->
    <mxCell id="llmsvc" value="&lt;b&gt;LLMService&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#F0FDF4;strokeColor=#059669;" vertex="1" parent="1">
      <mxGeometry x="740" y="360" width="360" height="130" as="geometry"/>
    </mxCell>
    <mxCell id="ll1" value="- client: AsyncOpenAI" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="llmsvc"><mxGeometry x="0" y="26" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="lld" value="──────────────────────────────" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=4;fontSize=10;" vertex="1" parent="llmsvc"><mxGeometry x="0" y="44" width="360" height="14" as="geometry"/></mxCell>
    <mxCell id="ll2" value="★ + generate(messages, timeout) → str  [max_tokens=8000]" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="llmsvc"><mxGeometry x="0" y="58" width="360" height="18" as="geometry"/></mxCell>
    <mxCell id="ll3" value="★ + stream(messages, timeout) → AsyncIterator  [max_tokens=8000]" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="llmsvc"><mxGeometry x="0" y="76" width="360" height="18" as="geometry"/></mxCell>

    <!-- ═══════ OCR STRATEGIES ═══════ -->
    <mxCell id="hdr3" value="Стратегии OCR (паттерн Strategy)" style="text;html=1;strokeColor=none;fillColor=#D97706;fontColor=#ffffff;align=center;fontSize=13;fontStyle=1;spacingLeft=4;" vertex="1" parent="1">
      <mxGeometry x="1120" y="20" width="400" height="34" as="geometry"/>
    </mxCell>

    <!-- OCRStrategy -->
    <mxCell id="ocrstrat" value="&lt;&lt;interface&gt;&gt;\n&lt;b&gt;OCRStrategy&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=40;fillColor=#FFF7ED;strokeColor=#D97706;" vertex="1" parent="1">
      <mxGeometry x="1200" y="64" width="280" height="100" as="geometry"/>
    </mxCell>
    <mxCell id="os1" value="+ recognize(image_path: str) → OcrResult" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="ocrstrat"><mxGeometry x="0" y="42" width="280" height="18" as="geometry"/></mxCell>
    <mxCell id="os2" value="+ get_name() → str" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="ocrstrat"><mxGeometry x="0" y="60" width="280" height="18" as="geometry"/></mxCell>

    <!-- PaddleOCRStrategy -->
    <mxCell id="paddle" value="&lt;b&gt;PaddleOCRStrategy&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#FFF7ED;strokeColor=#D97706;" vertex="1" parent="1">
      <mxGeometry x="1120" y="220" width="240" height="80" as="geometry"/>
    </mxCell>
    <mxCell id="pd1" value="+ recognize(image_path) → OcrResult" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="paddle"><mxGeometry x="0" y="26" width="240" height="18" as="geometry"/></mxCell>
    <mxCell id="pd2" value="+ get_name() → «Paddle»" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="paddle"><mxGeometry x="0" y="44" width="240" height="18" as="geometry"/></mxCell>

    <!-- TesseractStrategy -->
    <mxCell id="tess" value="&lt;b&gt;TesseractStrategy&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#FFF7ED;strokeColor=#D97706;" vertex="1" parent="1">
      <mxGeometry x="1380" y="220" width="240" height="80" as="geometry"/>
    </mxCell>
    <mxCell id="ts1" value="+ recognize(image_path) → OcrResult" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="tess"><mxGeometry x="0" y="26" width="240" height="18" as="geometry"/></mxCell>
    <mxCell id="ts2" value="+ get_name() → «Tesseract»" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="tess"><mxGeometry x="0" y="44" width="240" height="18" as="geometry"/></mxCell>

    <!-- HybridOCR -->
    <mxCell id="hybrid" value="&lt;b&gt;HybridOCR&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#FFF7ED;strokeColor=#D97706;" vertex="1" parent="1">
      <mxGeometry x="1200" y="360" width="280" height="120" as="geometry"/>
    </mxCell>
    <mxCell id="hy1" value="- strategies: list[OCRStrategy]  ◆" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="hybrid"><mxGeometry x="0" y="26" width="280" height="18" as="geometry"/></mxCell>
    <mxCell id="hyd" value="──────────────────────────────" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=4;fontSize=10;" vertex="1" parent="hybrid"><mxGeometry x="0" y="44" width="280" height="14" as="geometry"/></mxCell>
    <mxCell id="hy2" value="+ recognize(image_path) → OcrResult" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="hybrid"><mxGeometry x="0" y="58" width="280" height="18" as="geometry"/></mxCell>
    <mxCell id="hy3" value="  score = 0.7×conf + 0.3×text_quality" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=9;fontColor=#64748B;" vertex="1" parent="hybrid"><mxGeometry x="0" y="76" width="280" height="18" as="geometry"/></mxCell>

    <!-- ═══════ REPOSITORIES ═══════ -->
    <mxCell id="hdr4" value="Репозитории (паттерн Repository)" style="text;html=1;strokeColor=none;fillColor=#1E293B;fontColor=#ffffff;align=center;fontSize=13;fontStyle=1;spacingLeft=4;" vertex="1" parent="1">
      <mxGeometry x="20" y="560" width="700" height="34" as="geometry"/>
    </mxCell>

    <!-- CheckRepository -->
    <mxCell id="chrepo" value="&lt;b&gt;CheckRepository&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#F8FAFC;strokeColor=#1E293B;" vertex="1" parent="1">
      <mxGeometry x="20" y="604" width="330" height="150" as="geometry"/>
    </mxCell>
    <mxCell id="cr1" value="+ create(data: dict) → Check" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="chrepo"><mxGeometry x="0" y="26" width="330" height="18" as="geometry"/></mxCell>
    <mxCell id="cr2" value="+ get(check_id) → Check | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="chrepo"><mxGeometry x="0" y="44" width="330" height="18" as="geometry"/></mxCell>
    <mxCell id="cr3" value="+ get_by_user(user_id) → list[dict]" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="chrepo"><mxGeometry x="0" y="62" width="330" height="18" as="geometry"/></mxCell>
    <mxCell id="cr4" value="+ update(check_id, data) → Check | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="chrepo"><mxGeometry x="0" y="80" width="330" height="18" as="geometry"/></mxCell>
    <mxCell id="cr5" value="+ delete(check_id) → None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="chrepo"><mxGeometry x="0" y="98" width="330" height="18" as="geometry"/></mxCell>

    <!-- FunctionRepository -->
    <mxCell id="fnrepo" value="&lt;b&gt;FunctionRepository&lt;/b&gt;" style="swimlane;fontStyle=1;align=center;startSize=26;fillColor=#F8FAFC;strokeColor=#1E293B;" vertex="1" parent="1">
      <mxGeometry x="370" y="604" width="350" height="310" as="geometry"/>
    </mxCell>
    <mxCell id="fr1" value="+ list(user_id) → list[Function]" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="26" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr2" value="+ get(function_id) → Function | None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="44" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr3" value="+ create(data, user_id) → Function" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="62" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr4" value="+ update(function_id, data) → Function" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="80" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr5" value="+ delete(function_id) → None" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="98" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr6" value="+ publish(function_id, change_note) → Function" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="116" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr7" value="+ republish(function_id, change_note) → Function" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="134" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr8" value="+ copy(function_id, user_id) → Function" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="152" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr9" value="+ get_versions(function_id) → list[FunctionVersion]" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="170" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr10" value="★ + create_version(function_id, change_note)" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="188" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr11" value="★ + delete_version(version_id) → None" style="text;strokeColor=none;fillColor=#EEF2FF;align=left;spacingLeft=8;fontSize=10;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="206" width="350" height="18" as="geometry"/></mxCell>
    <mxCell id="fr12" value="+ list_gallery(user_id, search) → list[dict]" style="text;strokeColor=none;fillColor=none;align=left;spacingLeft=8;fontSize=10;" vertex="1" parent="fnrepo"><mxGeometry x="0" y="224" width="350" height="18" as="geometry"/></mxCell>

    <!-- Связи классов -->
    <!-- implements -->
    <mxCell id="impl1" style="edgeStyle=orthogonalEdgeStyle;endArrow=block;endFill=0;dashed=1;" edge="1" source="paddle" target="ocrstrat" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <mxCell id="impl2" style="edgeStyle=orthogonalEdgeStyle;endArrow=block;endFill=0;dashed=1;" edge="1" source="tess" target="ocrstrat" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <!-- composition HybridOCR -> OCRStrategy -->
    <mxCell id="comp1" style="edgeStyle=orthogonalEdgeStyle;endArrow=diamond;endFill=1;" edge="1" source="hybrid" target="ocrstrat" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <!-- CheckService -> LLMService -->
    <mxCell id="dep1" style="edgeStyle=orthogonalEdgeStyle;endArrow=open;dashed=1;" edge="1" source="csvc" target="llmsvc" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <!-- FunctionVersion -> Function -->
    <mxCell id="fk1" value="function_id (FK)" style="edgeStyle=orthogonalEdgeStyle;endArrow=open;fontSize=9;" edge="1" source="fver" target="func" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>
    <!-- Check -> Function -->
    <mxCell id="fk2" value="function_id (FK)" style="edgeStyle=orthogonalEdgeStyle;endArrow=open;fontSize=9;" edge="1" source="check" target="func" parent="1"><mxGeometry relative="1" as="geometry"/></mxCell>

    <!-- Легенда -->
    <mxCell id="leg2" value="★ — новые или обновлённые атрибуты / методы" style="text;html=1;strokeColor=none;fillColor=#EEF2FF;align=left;verticalAlign=middle;spacingLeft=8;fontSize=12;fontColor=#4F46E5;fontStyle=1;" vertex="1" parent="1">
      <mxGeometry x="20" y="960" width="500" height="36" as="geometry"/>
    </mxCell>
  </root>
</mxGraphModel>'''


def make_xml():
    uc_path = os.path.join(OUT, 'use_case_diagram.xml')
    cl_path = os.path.join(OUT, 'class_diagram.xml')
    with open(uc_path, 'w', encoding='utf-8') as f:
        f.write(USE_CASE_XML)
    print(f'✓ {uc_path}')
    # fix: add html=1 to swimlane styles so <b>...</b> renders as bold
    # fix: change field markers from + to - (ER-diagram convention)
    # fix: pupil_id → student_id
    cl_xml = CLASS_XML
    cl_xml = cl_xml.replace('style="swimlane;', 'style="swimlane;html=1;')
    cl_xml = cl_xml.replace('value="+ pupil_id: str | None"', 'value="- student_id: str | None"')
    # change public (+) to private (-) for all domain model field rows
    import re
    cl_xml = re.sub(r'value="\+ ([a-z_★])', r'value="- \1', cl_xml)
    with open(cl_path, 'w', encoding='utf-8') as f:
        f.write(cl_xml)
    print(f'✓ {cl_path}')


if __name__ == '__main__':
    make_pptx()
    make_docx()
    make_xml()
    print('\nГотово! Все файлы в папке diploma_updates/')
    print('XML-файлы открывай на сайте draw.io (File → Import from → Device)')
