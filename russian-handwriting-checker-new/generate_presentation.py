"""Генератор презентации дипломного проекта «РусЯзык AI» — 20 слайдов"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Цветовая палитра ───────────────────────────────────────────────────────────
INDIGO      = RGBColor(0x4F, 0x46, 0xE5)
INDIGO_D    = RGBColor(0x38, 0x2E, 0xB0)
INDIGO_L    = RGBColor(0xE0, 0xE7, 0xFF)
INDIGO_LL   = RGBColor(0xF5, 0xF3, 0xFF)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK  = RGBColor(0x1E, 0x1B, 0x4B)
DARK        = RGBColor(0x1E, 0x29, 0x3B)
GRAY        = RGBColor(0x64, 0x74, 0x8B)
GRAY_L      = RGBColor(0xF1, 0xF5, 0xF9)
GREEN       = RGBColor(0x05, 0x96, 0x69)
GREEN_L     = RGBColor(0xDC, 0xFC, 0xE7)
AMBER       = RGBColor(0xD9, 0x77, 0x06)
AMBER_L     = RGBColor(0xFF, 0xF7, 0xED)
RED         = RGBColor(0xDC, 0x26, 0x26)
RED_L       = RGBColor(0xFE, 0xF2, 0xF2)
TEAL        = RGBColor(0x0F, 0x76, 0x6E)
TEAL_L      = RGBColor(0xCC, 0xFB, 0xF1)
VIOLET      = RGBColor(0x7C, 0x3A, 0xED)
VIOLET_L    = RGBColor(0xED, 0xE9, 0xFE)
SLIDE_BG    = RGBColor(0xF8, 0xF9, 0xFF)

W = Inches(13.33)
H = Inches(7.5)


# ── Базовые примитивы ──────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def box(slide, x, y, w, h, color, border=False, border_color=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border and border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def tb(slide, x, y, w, h, text, sz=16, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    b  = slide.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text = text
    r.font.size   = Pt(sz)
    r.font.bold   = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name   = "Calibri"
    return b

def mtb(slide, x, y, w, h, lines, sz=14, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, spacing=None):
    b  = slide.shapes.add_textbox(x, y, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(sz)
        r.font.bold  = bold
        r.font.color.rgb = color
        r.font.name  = "Calibri"
    return b

def hdr(slide, title, subtitle=None):
    """Стандартная шапка слайда."""
    box(slide, 0, 0, W, Inches(1.18), INDIGO)
    tb(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.6),
       title, sz=24, bold=True, color=WHITE)
    if subtitle:
        tb(slide, Inches(0.5), Inches(0.78), Inches(12.3), Inches(0.32),
           subtitle, sz=13, color=INDIGO_L)

def num(slide, n, total=20):
    tb(slide, Inches(12.2), Inches(7.15), Inches(1.0), Inches(0.28),
       f"{n} / {total}", sz=10, color=GRAY, align=PP_ALIGN.RIGHT)

def divider(slide, x, y, w, color=INDIGO):
    box(slide, x, y, w, Pt(2), color)

# ── Карточки ──────────────────────────────────────────────────────────────────

def card(slide, x, y, w, h, title, body, accent=INDIGO,
         bg=WHITE, title_sz=14, body_sz=12):
    box(slide, x, y, w, h, bg)
    box(slide, x, y, w, Pt(5), accent)
    tb(slide, x + Inches(0.15), y + Inches(0.13),
       w - Inches(0.25), Inches(0.45),
       title, sz=title_sz, bold=True, color=accent)
    tb(slide, x + Inches(0.15), y + Inches(0.6),
       w - Inches(0.25), h - Inches(0.65),
       body, sz=body_sz, color=GRAY)

def icon_card(slide, x, y, w, h, icon, title, body, accent=INDIGO, bg=WHITE):
    box(slide, x, y, w, h, bg)
    box(slide, x, y, w, Pt(5), accent)
    tb(slide, x + Inches(0.12), y + Inches(0.1),
       Inches(0.5), Inches(0.48), icon, sz=24, align=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.65), y + Inches(0.14),
       w - Inches(0.75), Inches(0.38),
       title, sz=13, bold=True, color=accent)
    tb(slide, x + Inches(0.15), y + Inches(0.58),
       w - Inches(0.25), h - Inches(0.65),
       body, sz=11, color=GRAY)

def step_box(slide, x, y, w, h, num_str, title, detail, color=INDIGO):
    box(slide, x, y, w, h, WHITE)
    box(slide, x, y, w, Pt(5), color)
    box(slide, x + Inches(0.12), y + Inches(0.14),
        Inches(0.38), Inches(0.38), color)
    tb(slide, x + Inches(0.12), y + Inches(0.14),
       Inches(0.38), Inches(0.38),
       num_str, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(slide, x + Inches(0.58), y + Inches(0.18),
       w - Inches(0.7), Inches(0.38),
       title, sz=12, bold=True, color=color)
    tb(slide, x + Inches(0.12), y + Inches(0.62),
       w - Inches(0.22), h - Inches(0.7),
       detail, sz=10, color=GRAY)

def arrow_h(slide, x, y):
    tb(slide, x, y, Inches(0.3), Inches(0.35),
       "→", sz=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

def arrow_v(slide, x, y):
    tb(slide, x, y, Inches(0.35), Inches(0.3),
       "↓", sz=18, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

def tag(slide, x, y, w, h, text, bg_c, text_c=WHITE, sz=12):
    box(slide, x, y, w, h, bg_c)
    tb(slide, x + Inches(0.1), y + Inches(0.04),
       w - Inches(0.12), h, text, sz=sz, color=text_c, bold=True)

def code_block(slide, x, y, w, h, code_text, sz=10):
    box(slide, x, y, w, h, NEAR_BLACK)
    tb(slide, x + Inches(0.2), y + Inches(0.1),
       w - Inches(0.3), h - Inches(0.15),
       code_text, sz=sz, color=RGBColor(0xA5, 0xB4, 0xFC))

# ══════════════════════════════════════════════════════════════════════════════
# СЛАЙДЫ
# ══════════════════════════════════════════════════════════════════════════════

def s01_title(prs):
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO)
    box(sl, 0, Inches(5.3), W, Inches(2.2), INDIGO_D)
    box(sl, 0, 0, Inches(0.08), H, WHITE)

    tb(sl, Inches(0.5), Inches(0.6), Inches(12.3), Inches(1.1),
       "РусЯзык AI", sz=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    tb(sl, Inches(0.5), Inches(1.75), Inches(11.5), Inches(0.55),
       "Веб-приложение для автоматической проверки письменных работ учеников",
       sz=20, color=INDIGO_L)
    tb(sl, Inches(0.5), Inches(2.32), Inches(11.5), Inches(0.45),
       "с использованием OCR и искусственного интеллекта",
       sz=20, color=INDIGO_L)

    divider(sl, Inches(0.5), Inches(3.0), Inches(10), WHITE)

    mtb(sl, Inches(0.5), Inches(3.2), Inches(8), Inches(1.2),
        ["Дипломный проект",
         "Специальность: Разработка программного обеспечения",
         "2026 год"],
        sz=15, color=RGBColor(0xC7, 0xD2, 0xFE), spacing=22)

    for i, (icon, lbl) in enumerate([("🖊", "OCR\nрукописей"),
                                      ("🤖", "ИИ\nпроверка"),
                                      ("📊", "ОГЭ / ЕГЭ\nкритерии"),
                                      ("👨‍🏫", "Управление\nучениками")]):
        x = Inches(0.5) + i * Inches(3.1)
        box(sl, x, Inches(5.5), Inches(2.8), Inches(1.6), INDIGO_D)
        tb(sl, x + Inches(0.15), Inches(5.6), Inches(0.55), Inches(0.55),
           icon, sz=28)
        tb(sl, x + Inches(0.75), Inches(5.63), Inches(1.9), Inches(1.3),
           lbl, sz=13, color=INDIGO_L)

    num(sl, 1)

def s02_problem(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Постановка проблемы")

    box(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.95), INDIGO_L)
    tb(sl, Inches(0.75), Inches(1.42), Inches(11.8), Inches(0.72),
       "Учитель русского языка тратит 10–15 минут на проверку одной работы.\n"
       "При классе из 30 учеников — это 5–7 часов только на одну волну работ.",
       sz=17, bold=True, color=INDIGO)

    problems = [
        ("📷", "Прочитать\nрукописный\nпочерк",        INDIGO),
        ("✏️",  "Найти каждую\nошибку и\nобъяснить её", VIOLET),
        ("📊", "Оценить по\nкритериям\nОГЭ / ЕГЭ",     TEAL),
        ("🗂",  "Сохранить\nрезультаты\nдля отчёта",    AMBER),
    ]
    cw = Inches(2.9)
    ch = Inches(2.4)
    gx = Inches(0.3)
    y0 = Inches(2.45)
    for i, (icon, text, col) in enumerate(problems):
        x = Inches(0.5) + i * (cw + gx)
        box(sl, x, y0, cw, ch, WHITE)
        box(sl, x, y0, cw, Pt(6), col)
        tb(sl, x + Inches(0.1), y0 + Inches(0.15),
           cw - Inches(0.2), Inches(0.6),
           icon, sz=30, align=PP_ALIGN.CENTER)
        tb(sl, x + Inches(0.1), y0 + Inches(0.8),
           cw - Inches(0.2), Inches(1.4),
           text, sz=14, color=col, bold=True, align=PP_ALIGN.CENTER)

    box(sl, 0, Inches(5.9), W, Inches(1.6), NEAR_BLACK)
    tb(sl, Inches(0.7), Inches(6.05), Inches(11.9), Inches(1.3),
       "Цель проекта: автоматизировать рутинную часть проверки, "
       "чтобы учитель мог\nсосредоточиться на живом взаимодействии с учеником",
       sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    num(sl, 2)

def s03_goals(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Цель и задачи проекта")

    box(sl, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.82), INDIGO_L)
    box(sl, Inches(0.5), Inches(1.3), Pt(6), Inches(0.82), INDIGO)
    tb(sl, Inches(0.75), Inches(1.4), Inches(11.8), Inches(0.65),
       "ЦЕЛЬ: разработать веб-приложение для автоматической проверки рукописных "
       "и печатных работ учеников по русскому языку",
       sz=16, bold=True, color=INDIGO)

    tasks = [
        ("Распознавание текста",   "Рукописные фото → HybridOCR\nPDF/DOCX/PPTX → MarkItDown"),
        ("Проверка по критериям",  "ОГЭ, ЕГЭ, итоговое сочинение,\nорфография, свободная проверка"),
        ("Стриминг в реальном времени", "Ответ ИИ появляется\nпо мере генерации"),
        ("История и статистика",   "Хранение всех проверок,\nпоиск и фильтрация"),
        ("Управление учениками",   "База учеников, группы (классы),\nпривязка работ к ученику"),
        ("Галерея функций",        "Создать, опубликовать,\nскопировать, версионировать"),
    ]
    cw = Inches(3.9)
    ch = Inches(1.52)
    gx = Inches(0.22)
    gy = Inches(0.18)
    pos = [(Inches(0.5) + i*(cw+gx), Inches(2.35) + j*(ch+gy))
           for j in range(2) for i in range(3)]
    colors = [INDIGO, VIOLET, TEAL, AMBER, GREEN, RED]

    for (x, y), (title, body), col in zip(pos, tasks, colors):
        box(sl, x, y, cw, ch, WHITE)
        box(sl, x, y, cw, Pt(6), col)
        box(sl, x, y, Inches(0.06), ch, col)
        tb(sl, x + Inches(0.15), y + Inches(0.1),
           cw - Inches(0.22), Inches(0.45),
           title, sz=14, bold=True, color=col)
        tb(sl, x + Inches(0.15), y + Inches(0.6),
           cw - Inches(0.22), Inches(0.82),
           body, sz=12, color=GRAY)

    num(sl, 3)

def s04_competitors(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Анализ конкурентов")

    cols_x = [Inches(0.5), Inches(3.0), Inches(5.6), Inches(8.2), Inches(10.5)]
    cols_w = [Inches(2.4), Inches(2.5), Inches(2.5), Inches(2.2), Inches(2.7)]
    hdrs_  = ["Решение", "OCR\nрукописей", "Критерии\nОГЭ/ЕГЭ",
              "История\nи ученики", "Свои\nфункции"]

    y_h = Inches(1.3)
    for cw, cx, ht in zip(cols_w, cols_x, hdrs_):
        box(sl, cx, y_h, cw - Inches(0.04), Inches(0.7), INDIGO)
        tb(sl, cx + Inches(0.08), y_h + Inches(0.05),
           cw - Inches(0.15), Inches(0.62),
           ht, sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Яндекс.Спеллер",   "✗", "✗", "✗", "✗"),
        ("LanguageTool",      "✗", "✗", "✗", "частично"),
        ("ChatGPT (ручной)", "✗", "частично", "✗", "✗"),
        ("Google Vision API", "✓", "✗", "✗", "✗"),
        ("Tesseract (авт.)", "частично", "✗", "✗", "✗"),
        ("РусЯзык AI ★",    "✓", "✓", "✓", "✓"),
    ]

    for r, row in enumerate(rows):
        y = Inches(2.05) + r * Inches(0.82)
        is_us = (r == len(rows) - 1)
        rb = INDIGO_L if is_us else (WHITE if r % 2 == 0 else GRAY_L)
        for cw, cx, cell in zip(cols_w, cols_x, row):
            box(sl, cx, y, cw - Inches(0.04), Inches(0.76), rb)
            c = INDIGO if cell == row[0] and is_us else \
                GREEN if cell == "✓" else \
                RED if cell == "✗" else \
                AMBER if cell == "частично" else \
                (INDIGO if is_us else DARK)
            tb(sl, cx + Inches(0.08), y + Inches(0.16),
               cw - Inches(0.15), Inches(0.48),
               cell, sz=13 if cell == row[0] else 17,
               bold=is_us, color=c, align=PP_ALIGN.CENTER)

    num(sl, 4)

def s05_features(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Функциональность приложения",
        subtitle="Полный набор инструментов учителя в одном месте")

    feats = [
        ("📤", "Загрузка файлов",         "Фото, PDF, DOCX, PPTX, TXT —\nлюбой формат"),
        ("✍️",  "Редактор текста",          "Tiptap: отредактировать\nраспознанный текст"),
        ("🤖", "ИИ-проверка",              "Стриминг результата —\nпоявляется в реальном времени"),
        ("💾", "Сохранение работ",         "Привязка к ученику, папке,\nдате и оценке"),
        ("👨‍🎓", "База учеников",            "Список, группы, статистика\nпо каждому ученику"),
        ("⚙️",  "Функции проверки",         "Создать, опубликовать\nв галерею, копировать"),
        ("📁",  "Папки",                   "Организовать работы\nпо темам или классам"),
        ("🖨️",  "Печать отчёта",           "HTML-отчёт прямо\nиз браузера"),
    ]

    cw = Inches(2.95)
    ch = Inches(1.55)
    gx = Inches(0.22)
    gy = Inches(0.18)
    colors = [INDIGO, TEAL, VIOLET, GREEN, AMBER, RED, INDIGO, TEAL]

    for i, ((icon, title, body), col) in enumerate(zip(feats, colors)):
        col_i = i % 4
        row_i = i // 4
        x = Inches(0.5) + col_i * (cw + gx)
        y = Inches(1.5) + row_i * (ch + gy)
        box(sl, x, y, cw, ch, WHITE)
        box(sl, x, y, cw, Pt(6), col)
        tb(sl, x + Inches(0.12), y + Inches(0.12),
           Inches(0.48), Inches(0.48), icon, sz=24)
        tb(sl, x + Inches(0.65), y + Inches(0.15),
           cw - Inches(0.75), Inches(0.38),
           title, sz=13, bold=True, color=col)
        tb(sl, x + Inches(0.15), y + Inches(0.62),
           cw - Inches(0.25), Inches(0.82),
           body, sz=11, color=GRAY)

    num(sl, 5)

def s06_check_types(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Типы проверки: встроенные функции",
        subtitle="5 функций готовы сразу после запуска системы (сеедирование БД)")

    types = [
        (INDIGO,  "Орфография\nи пунктуация",
         "Ошибки с объяснением\nОценка 0–5\nОбщий комментарий"),
        (GREEN,   "Сочинение\nОГЭ 9.3",
         "8 критериев:\nСК1–СК4 (содержание)\nГК1–ГК4 (грамотность)"),
        (AMBER,   "Сочинение\nЕГЭ",
         "12 критериев К1–К12\nпо нормам ФИПИ\nСумма баллов"),
        (VIOLET,  "Итоговое\nсочинение",
         "5 критериев К1–К5\nЗачёт / незачёт\nДопуск к ЕГЭ"),
        (TEAL,    "Свободная\nпроверка",
         "Учитель сам\nформулирует задание\nИИ его выполняет"),
    ]

    cw = Inches(2.32)
    ch = Inches(4.7)
    gx = Inches(0.25)
    y0 = Inches(1.5)

    for i, (col, title, desc) in enumerate(types):
        x = Inches(0.5) + i * (cw + gx)
        box(sl, x, y0, cw, ch, WHITE)
        box(sl, x, y0, cw, Pt(8), col)
        box(sl, x, y0, cw, Inches(1.5), col)
        tb(sl, x + Inches(0.12), y0 + Inches(0.2),
           cw - Inches(0.2), Inches(1.15),
           title, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x + Inches(0.12), y0 + Inches(1.65),
           cw - Inches(0.2), ch - Inches(1.8),
           desc, sz=12, color=GRAY)

    box(sl, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.75), INDIGO_L)
    tb(sl, Inches(0.75), Inches(6.5), Inches(11.8), Inches(0.55),
       "✚  Пользователь может создать собственную функцию, "
       "опубликовать в галерею и поделиться с другими учителями",
       sz=14, bold=True, color=INDIGO)

    num(sl, 6)

def s07_architecture(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Архитектура системы: трёхуровневая",
        subtitle="Разделение ответственности: Routes → Services → Repositories")

    # ── Три колонки архитектуры ────────────────────────────────────────────────
    # Ширина колонки = 2.9", зазор = 0.2", старт = 0.4"
    # Col 0: 0.40 – 3.30   Col 1: 3.50 – 6.40   Col 2: 6.60 – 9.50
    # Инфраструктура: 9.75 – 13.10  (не пересекается!)
    layers = [
        (INDIGO, "ROUTES\n(Маршруты)",
         ["• Принимает HTTP-запрос",
          "• Валидирует (Pydantic)",
          "• Делегирует сервису",
          "• Возвращает JSON / SSE",
          "",
          "auth  /  check  /  upload",
          "functions  /  groups",
          "pupils  /  folders"]),
        (TEAL,   "SERVICES\n(Бизнес-логика)",
         ["• CheckService",
          "  промпт + парсинг JSON",
          "• LLMService — стриминг",
          "• DocumentConverter",
          "  OCR или MarkItDown",
          "• HybridOCR",
          "  PaddleOCR + Tesseract"]),
        (VIOLET, "REPOSITORIES\n(Доступ к данным)",
         ["• CheckRepository",
          "• FunctionRepository",
          "• GroupRepository",
          "• PupilRepository",
          "• FolderRepository",
          "",
          "SQL скрыт от логики"]),
    ]

    cw = Inches(2.9)   # ширина одной колонки
    ch = Inches(5.55)  # высота колонки
    gx = Inches(0.2)   # зазор между колонками
    y0 = Inches(1.42)
    x0 = Inches(0.4)

    for i, (col, title, items) in enumerate(layers):
        x = x0 + i * (cw + gx)
        box(sl, x, y0, cw, ch, WHITE)
        box(sl, x, y0, cw, Inches(0.72), col)   # цветной заголовок
        box(sl, x, y0, Pt(4), ch, col)           # левая полоска
        tb(sl, x + Inches(0.14), y0 + Inches(0.1),
           cw - Inches(0.22), Inches(0.58),
           title, sz=13, bold=True, color=WHITE)
        y_i = y0 + Inches(0.85)
        for item in items:
            is_bullet = item.startswith("•")
            is_sub    = item.startswith("  ")
            tb(sl, x + Inches(0.15), y_i,
               cw - Inches(0.25), Inches(0.52),
               item,
               sz=11,
               color=DARK if is_bullet else (GRAY if is_sub else RGBColor(0x94, 0xA3, 0xB8)))
            y_i += Inches(0.53)

        # Стрелка между колонками
        if i < 2:
            ax = x + cw + Inches(0.02)
            tb(sl, ax, y0 + Inches(2.5), Inches(0.22), Inches(0.42),
               "→", sz=20, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)

    # ── Инфраструктура (правая панель, x=9.75) ─────────────────────────────────
    ix = Inches(9.75)
    iw = Inches(3.35)
    ih_title = Inches(0.45)

    box(sl, ix, y0, iw, ih_title, INDIGO_D)
    tb(sl, ix + Inches(0.12), y0 + Inches(0.05),
       iw - Inches(0.15), ih_title - Inches(0.06),
       "ИНФРАСТРУКТУРА", sz=12, bold=True, color=WHITE)

    infra = [
        (INDIGO,  INDIGO_L,  "🔐  Supabase Auth",    "JWT-авторизация"),
        (GREEN,   GREEN_L,   "🗄  PostgreSQL",        "Supabase + SQLite (dev)"),
        (VIOLET,  VIOLET_L,  "🤖  Ollama qwen2.5:7b","Локальный LLM"),
        (AMBER,   AMBER_L,   "🐳  Docker",            "docker-compose up"),
    ]

    y_inf = y0 + ih_title + Inches(0.1)
    for acc, bg_c, title_i, desc_i in infra:
        box(sl, ix, y_inf, iw, Inches(1.18), bg_c)
        box(sl, ix, y_inf, Pt(5), Inches(1.18), acc)
        tb(sl, ix + Inches(0.14), y_inf + Inches(0.1),
           iw - Inches(0.2), Inches(0.42),
           title_i, sz=13, bold=True, color=acc)
        tb(sl, ix + Inches(0.14), y_inf + Inches(0.55),
           iw - Inches(0.2), Inches(0.52),
           desc_i, sz=11, color=GRAY)
        y_inf += Inches(1.25)

    num(sl, 7)

def s08_ocr_pipeline(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Обработка файлов: алгоритм DocumentConverter",
        subtitle="Два пути + fallback — единый интерфейс для любого формата")

    # Входной файл
    box(sl, Inches(0.5), Inches(1.5), Inches(2.4), Inches(0.82), INDIGO)
    tb(sl, Inches(0.6), Inches(1.62), Inches(2.2), Inches(0.62),
       "📂  Загруженный\nфайл", sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    arrow_v(sl, Inches(1.52), Inches(2.35))

    # DocumentConverter
    box(sl, Inches(0.5), Inches(2.72), Inches(2.4), Inches(0.72), INDIGO_D)
    tb(sl, Inches(0.6), Inches(2.82), Inches(2.2), Inches(0.55),
       "DocumentConverter", sz=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Разделение
    tb(sl, Inches(2.95), Inches(2.88), Inches(0.3), Inches(0.42),
       "→", sz=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

    # ПУТЬ 1 — изображение
    y1 = Inches(1.5)
    box(sl, Inches(3.4), y1, Inches(3.8), Inches(0.72), INDIGO_L)
    tb(sl, Inches(3.55), y1 + Inches(0.1), Inches(3.6), Inches(0.55),
       "🖼  Изображение: .jpg  .png  .bmp  .tiff",
       sz=12, color=INDIGO, bold=True)

    arrow_v(sl, Inches(5.15), y1 + Inches(0.74))

    box(sl, Inches(3.4), y1 + Inches(1.12), Inches(3.8), Inches(0.72),
        RGBColor(0xE0, 0xF2, 0xFE))
    tb(sl, Inches(3.55), y1 + Inches(1.22), Inches(3.6), Inches(0.55),
       "HybridOCR", sz=14, bold=True, color=RGBColor(0x07, 0x52, 0x85))

    arrow_v(sl, Inches(5.15), y1 + Inches(1.87))

    # Два движка
    for idx, (lbl, col) in enumerate([("PaddleOCR\n2.8.1", RGBColor(0x07, 0x52, 0x85)),
                                       ("Tesseract\n5.x",   RGBColor(0x16, 0x45, 0x7A))]):
        x = Inches(3.4) + idx * Inches(2.0)
        box(sl, x, y1 + Inches(2.28), Inches(1.85), Inches(0.85),
            RGBColor(0xBA, 0xE6, 0xFD))
        tb(sl, x + Inches(0.1), y1 + Inches(2.36), Inches(1.65), Inches(0.72),
           lbl, sz=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    arrow_v(sl, Inches(5.15), y1 + Inches(3.15))
    box(sl, Inches(3.4), y1 + Inches(3.55), Inches(3.8), Inches(0.62),
        GREEN_L)
    tb(sl, Inches(3.55), y1 + Inches(3.65), Inches(3.6), Inches(0.48),
       "score = 0.7×conf + 0.3×text_quality → лучший",
       sz=10, color=GREEN, italic=True)

    # ПУТЬ 2 — документы
    y2 = Inches(1.5)
    box(sl, Inches(7.55), y2, Inches(4.2), Inches(0.72), INDIGO_L)
    tb(sl, Inches(7.7), y2 + Inches(0.1), Inches(4.0), Inches(0.55),
       "📄  Документ: .pdf  .docx  .pptx  .txt  .md",
       sz=12, color=INDIGO, bold=True)

    arrow_v(sl, Inches(9.5), y2 + Inches(0.74))

    box(sl, Inches(7.55), y2 + Inches(1.12), Inches(4.2), Inches(0.72), GREEN_L)
    tb(sl, Inches(7.7), y2 + Inches(1.22), Inches(4.0), Inches(0.55),
       "MarkItDown (Microsoft)\nПDF/DOCX/PPTX → Markdown-текст",
       sz=11, color=GREEN)

    arrow_v(sl, Inches(9.5), y2 + Inches(1.87))

    box(sl, Inches(7.55), y2 + Inches(2.28), Inches(4.2), Inches(0.72), GREEN_L)
    box(sl, Inches(7.55), y2 + Inches(2.28), Inches(4.2), Pt(5), GREEN)
    tb(sl, Inches(7.7), y2 + Inches(2.38), Inches(4.0), Inches(0.55),
       "✓  Текст извлечён", sz=13, bold=True, color=GREEN)

    # fallback
    tb(sl, Inches(9.8), y2 + Inches(2.5), Inches(0.4), Inches(0.42),
       "✗", sz=18, bold=True, color=RED, align=PP_ALIGN.CENTER)
    arrow_v(sl, Inches(7.65), y2 + Inches(3.08))
    box(sl, Inches(7.55), y2 + Inches(3.48), Inches(4.2), Inches(0.72), AMBER_L)
    box(sl, Inches(7.55), y2 + Inches(3.48), Inches(4.2), Pt(5), AMBER)
    tb(sl, Inches(7.7), y2 + Inches(3.58), Inches(4.0), Inches(0.55),
       "Fallback → HybridOCR", sz=13, bold=True, color=AMBER)

    # Итог
    box(sl, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.72), GREEN_L)
    box(sl, Inches(0.5), Inches(5.85), Pt(5), Inches(0.72), GREEN)
    tb(sl, Inches(0.75), Inches(5.97), Inches(11.8), Inches(0.52),
       "✓  Извлечённый текст → передаётся в языковую модель (LLM) для проверки",
       sz=15, bold=True, color=GREEN)

    num(sl, 8)

def s09_streaming(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Стриминг ответа ИИ в реальном времени",
        subtitle="Server-Sent Events (SSE) — пользователь видит каждый токен сразу")

    steps = [
        (INDIGO,  "1",  "Нажата\nкнопка",
         "POST /api/check/stream\n{text, function_id}"),
        (TEAL,    "2",  "CheckService\nформирует промпт",
         "system_prompt + user_template\nс подставленным текстом"),
        (VIOLET,  "3",  "LLMService\nзапрашивает LLM",
         "openai.AsyncOpenAI\nstream=True, timeout=180s"),
        (AMBER,   "4",  "Токены →\nServer-Sent Events",
         'data: {"type":"chunk",\n "text":"..."}'),
        (GREEN,   "5",  "Браузер\nчитает поток",
         "fetch() + ReadableStream\n+ TextDecoder"),
    ]

    cw = Inches(2.32)
    ch = Inches(3.0)
    gx = Inches(0.22)
    y0 = Inches(1.5)

    for i, (col, n, title, detail) in enumerate(steps):
        x = Inches(0.4) + i * (cw + gx)
        step_box(sl, x, y0, cw, ch, n, title, detail, col)
        if i < 4:
            arrow_h(sl, x + cw + Inches(0.0), y0 + Inches(1.2))

    # Итоговый эффект
    box(sl, Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.75), INDIGO_L)
    tb(sl, Inches(0.75), Inches(4.77), Inches(11.8), Inches(0.55),
       "Без стриминга: пустой экран 30–60 секунд → весь текст сразу\n"
       "Со стримингом: текст появляется сразу, как при живой печати",
       sz=13, bold=True, color=INDIGO)

    # Код
    code_block(sl, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.65),
       "async def stream_check(text, function_id) -> AsyncIterator[dict]:\n"
       "    messages = await self._build_messages(text, function_id)   # system + user\n"
       "    full_text = ''\n"
       "    async for chunk in self.llm.stream(messages):              # кусочки от LLM\n"
       "        full_text += chunk\n"
       "        yield {'type': 'chunk', 'text': chunk}                 # → клиенту сразу\n"
       "    result = self._build_result(text, self._parse(full_text))\n"
       "    yield {'type': 'done', 'result': result}                   # финальный JSON")

    num(sl, 9)

def s10_upload_algo(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Алгоритм загрузки и обработки файла",
        subtitle="BackgroundTasks FastAPI + polling — клиент не ждёт OCR блокируя соединение")

    steps6 = [
        (INDIGO,  "1", "Пользователь\nвыбирает файл",
         "Drag & drop или диалог\nвыбора файла"),
        (INDIGO,  "2", "POST /api/upload/",
         "Файл сохраняется во temp/\nСоздаётся task_id\nОтвет: {task_id, status:'processing'}"),
        (TEAL,    "3", "Клиент опрашивает\nGET /upload/{task_id}",
         "Интервал 1.5 сек.\nПока status == 'processing'"),
        (VIOLET,  "4", "DocumentConverter\nв фоне",
         "Определяет тип файла\nHybridOCR или MarkItDown\nВыполняется параллельно"),
        (GREEN,   "5", "status → 'done'\ntemp-файл удалён",
         "Результат в tasks[task_id]\nВременный файл очищен"),
        (AMBER,   "6", "Текст вставляется\nв редактор",
         "editor.commands.\nsetContent(text)\nГотово к проверке"),
    ]

    cw = Inches(2.05)
    ch = Inches(2.85)
    gx = Inches(0.18)
    y0 = Inches(1.45)

    for i, (col, n, title, detail) in enumerate(steps6):
        x = Inches(0.38) + i * (cw + gx)
        step_box(sl, x, y0, cw, ch, n, title, detail, col)
        if i < 5:
            arrow_h(sl, x + cw + Inches(0.0), y0 + Inches(1.2))

    # Зачем polling?
    box(sl, Inches(0.5), Inches(4.45), Inches(12.3), Inches(1.05), INDIGO_LL)
    box(sl, Inches(0.5), Inches(4.45), Pt(5), Inches(1.05), INDIGO)
    mtb(sl, Inches(0.75), Inches(4.55), Inches(11.8), Inches(0.9),
        ["Зачем polling, а не держать соединение открытым?",
         "OCR (особенно PaddleOCR) занимает 5–30 секунд. "
         "Браузер закрывает HTTP-соединение по таймауту.\n"
         "Polling: быстрый ответ сразу → периодическая проверка → результат когда готов."],
        sz=12, color=DARK)

    # Код POST /upload
    code_block(sl, Inches(0.5), Inches(5.65), Inches(12.3), Inches(1.6),
        "@router.post('/')\n"
        "async def upload(file: UploadFile, background_tasks: BackgroundTasks):\n"
        "    task_id = str(uuid4())\n"
        "    path = TEMP_DIR / f'{task_id}_{file.filename}'\n"
        "    # сохраняем файл\n"
        "    tasks[task_id] = {'status': 'processing'}\n"
        "    background_tasks.add_task(ocr_service.process_file, str(path), task_id)\n"
        "    return {'task_id': task_id, 'status': 'processing'}")

    num(sl, 10)

def s11_check_service(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Как работает проверка: CheckService",
        subtitle="Routes → CheckService → LLMService → parse JSON → highlight errors")

    # Схема потока
    flow = [
        (INDIGO,  "POST\n/check/stream",         Inches(0.4)),
        (TEAL,    "FunctionRepository\n.get(id)", Inches(2.7)),
        (VIOLET,  "_build_messages()\nsystem + user", Inches(5.0)),
        (AMBER,   "LLMService\n.stream()",        Inches(7.3)),
        (GREEN,   "_parse()\nJSON → dict",        Inches(9.6)),
    ]
    y0 = Inches(1.45)
    for col, lbl, x in flow:
        box(sl, x, y0, Inches(2.1), Inches(0.95), col)
        tb(sl, x + Inches(0.1), y0 + Inches(0.1),
           Inches(1.9), Inches(0.78), lbl, sz=12, bold=True,
           color=WHITE, align=PP_ALIGN.CENTER)
    for i in range(4):
        arrow_h(sl, Inches(0.4) + (i + 1) * Inches(2.3) - Inches(0.2),
                y0 + Inches(0.35))

    # Поля результата
    box(sl, Inches(0.5), Inches(2.65), Inches(12.3), Inches(0.38), DARK)
    tb(sl, Inches(0.75), Inches(2.72), Inches(11.8), Inches(0.28),
       "_build_result() — что входит в итоговый ответ:", sz=13, bold=True, color=WHITE)

    fields = [
        ("corrected_text",   "Исправленный текст",             INDIGO_L,   INDIGO),
        ("errors",           "Список ошибок (JSON-массив)",    GREEN_L,    GREEN),
        ("score",            "Числовая оценка",                AMBER_L,    AMBER),
        ("criteria",         "Критерии (ОГЭ/ЕГЭ)",            VIOLET_L,   VIOLET),
        ("html_highlighted", "Текст с <span> вокруг ошибок",  TEAL_L,     TEAL),
        ("comment",          "Общий комментарий",              GRAY_L,     GRAY),
    ]
    cw = Inches(3.9)
    ch = Inches(0.92)
    gx = Inches(0.22)
    gy = Inches(0.14)
    for i, (key, desc, bg_c, acc) in enumerate(fields):
        col_i = i % 3
        row_i = i // 3
        x = Inches(0.5) + col_i * (cw + gx)
        y = Inches(3.15) + row_i * (ch + gy)
        box(sl, x, y, cw, ch, bg_c)
        box(sl, x, y, cw, Pt(4), acc)
        tb(sl, x + Inches(0.15), y + Inches(0.08),
           cw - Inches(0.2), Inches(0.4),
           key, sz=12, bold=True, color=acc)
        tb(sl, x + Inches(0.15), y + Inches(0.5),
           cw - Inches(0.2), Inches(0.38),
           desc, sz=11, color=GRAY)

    # Highlight code
    code_block(sl, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.82),
        "def _highlight(self, text: str, errors: list) -> str:\n"
        "    result = text\n"
        "    for e in errors:\n"
        "        original = e.get('original', '')          # слово с ошибкой из JSON\n"
        "        result = re.sub(\n"
        "            re.escape(original),                  # экранируем спецсимволы\n"
        "            f'<span class=\"error-highlight\">{original}</span>',\n"
        "            result, count=1)                      # только первое вхождение\n"
        "    return result")

    num(sl, 11)

def s12_functions(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Функции проверки: создание, публикация, галерея",
        subtitle="Каждая функция — это system_prompt + user_template с плейсхолдером {text}")

    # Жизненный цикл функции
    lifecycle = [
        (INDIGO,  "Создание",    "Учитель пишет system_prompt\n(инструкции для ИИ)\nи user_template\nс {text}"),
        (TEAL,    "Использование","Выбирается в CheckPanel\nТекст подставляется\nвместо {text}\nОтправляется в LLM"),
        (VIOLET,  "Публикация",  "Флаг is_published = True\nФункция попадает\nв галерею\nДругие видят её"),
        (AMBER,   "Версионирование","Редактирование уже\nопубликованной →\nавтоматически создаётся\nновая версия в function_versions"),
        (GREEN,   "Копирование", "Другой учитель\nнажимает «Скопировать»\nPOST /functions/{id}/copy\nКопия в его библиотеке"),
    ]

    cw = Inches(2.32)
    ch = Inches(3.5)
    gx = Inches(0.25)
    y0 = Inches(1.45)

    for i, (col, title, body) in enumerate(lifecycle):
        x = Inches(0.5) + i * (cw + gx)
        box(sl, x, y0, cw, ch, WHITE)
        box(sl, x, y0, cw, Inches(0.75), col)
        tb(sl, x + Inches(0.12), y0 + Inches(0.1),
           cw - Inches(0.2), Inches(0.58),
           title, sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(sl, x + Inches(0.12), y0 + Inches(0.88),
           cw - Inches(0.2), ch - Inches(0.95),
           body, sz=12, color=GRAY)
        if i < 4:
            arrow_h(sl, x + cw + Inches(0.0), y0 + Inches(1.55))

    # Защита стандартных
    box(sl, Inches(0.5), Inches(5.1), Inches(5.8), Inches(0.82), AMBER_L)
    box(sl, Inches(0.5), Inches(5.1), Pt(5), Inches(0.82), AMBER)
    tb(sl, Inches(0.7), Inches(5.2), Inches(5.5), Inches(0.65),
       "⚠  is_default = True → нельзя редактировать и удалять\n"
       "(5 стандартных функций защищены от изменений)", sz=12, color=AMBER)

    # Пример промпта
    code_block(sl, Inches(6.5), Inches(5.1), Inches(6.7), Inches(2.15),
        '# Пример: Проверка орфографии\n'
        'system_prompt: "Ты — учитель русского языка.\n'
        '  Верни JSON: {corrected, errors:[{original,\n'
        '  corrected, type, comment}], score, comment}"\n\n'
        'user_template: "Проверь следующий текст:\\n\\n{text}"')

    num(sl, 12)

def s13_frontend(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Клиентская часть: компоненты и маршруты",
        subtitle="React 19 + TypeScript + Tiptap + Tailwind CSS")

    # Страницы / роуты
    pages = [
        ("/",          "CheckPage",     "Главная страница проверки"),
        ("/history",   "HistoryPage",   "История всех проверок"),
        ("/students",  "StudentsPage",  "Ученики и группы"),
        ("/functions", "FunctionsPage", "Мои функции + галерея"),
        ("/profile",   "ProfilePage",   "Профиль и статистика"),
    ]
    box(sl, Inches(0.5), Inches(1.42), Inches(3.9), Inches(5.8), INDIGO_LL)
    box(sl, Inches(0.5), Inches(1.42), Pt(5), Inches(5.8), INDIGO)
    tb(sl, Inches(0.68), Inches(1.5), Inches(3.6), Inches(0.38),
       "React Router DOM — маршруты", sz=13, bold=True, color=INDIGO)
    for i, (route, comp, desc) in enumerate(pages):
        y = Inches(1.98) + i * Inches(1.0)
        box(sl, Inches(0.6), y, Inches(3.7), Inches(0.88), WHITE)
        tb(sl, Inches(0.75), y + Inches(0.05), Inches(3.4), Inches(0.32),
           route, sz=12, bold=True, color=INDIGO)
        tb(sl, Inches(0.75), y + Inches(0.38), Inches(3.4), Inches(0.32),
           f"{comp}  —  {desc}", sz=11, color=GRAY)

    # Компоненты CheckPage
    comps = [
        (INDIGO,  "UploadForm",       "Загрузка файла\ndraq & drop"),
        (TEAL,    "TextEditor",       "TipTap редактор\nOCR-текста"),
        (VIOLET,  "CheckPanel",       "Выбор функции\nи запуск проверки"),
        (AMBER,   "StreamingPreview", "Стриминг ответа\nв реальном времени"),
        (GREEN,   "ResultPanel",      "Результат:\nошибки, оценка,\nсохранение"),
    ]
    box(sl, Inches(4.65), Inches(1.42), Inches(8.2), Inches(0.38), TEAL)
    tb(sl, Inches(4.8), Inches(1.47), Inches(8.0), Inches(0.3),
       "CheckPage — компоненты главной страницы", sz=12, bold=True, color=WHITE)

    cw = Inches(1.52)
    ch = Inches(1.65)
    gx = Inches(0.12)
    for i, (col, name, desc) in enumerate(comps):
        x = Inches(4.65) + i * (cw + gx)
        box(sl, x, Inches(1.9), cw, ch, WHITE)
        box(sl, x, Inches(1.9), cw, Pt(5), col)
        tb(sl, x + Inches(0.08), Inches(1.97),
           cw - Inches(0.12), Inches(0.42),
           name, sz=10, bold=True, color=col)
        tb(sl, x + Inches(0.08), Inches(2.42),
           cw - Inches(0.12), Inches(1.0),
           desc, sz=10, color=GRAY)

    # AuthContext + Axios
    box(sl, Inches(4.65), Inches(3.72), Inches(8.2), Inches(1.55), INDIGO_LL)
    box(sl, Inches(4.65), Inches(3.72), Pt(5), Inches(1.55), INDIGO)
    mtb(sl, Inches(4.85), Inches(3.8), Inches(7.9), Inches(1.4),
        ["AuthContext (React Context API)",
         "JWT-токен хранится в localStorage.\n"
         "Axios-перехватчик автоматически добавляет\n"
         "Authorization: Bearer <token> к каждому запросу.",
         "При 401 — выход из системы."],
        sz=12, color=DARK)

    # TypeScript типы
    box(sl, Inches(4.65), Inches(5.42), Inches(8.2), Inches(1.82), NEAR_BLACK)
    tb(sl, Inches(4.82), Inches(5.52), Inches(8.0), Inches(1.65),
       "// types.ts — все API-ответы строго типизированы\n"
       "interface CheckRecord { id: string; score: number | null;\n"
       "  errors: ...; corrected_text?: string; ... }\n"
       "interface Fn { id: string; name: string; system_prompt?: string; ... }\n"
       "interface Pupil { id: string; name: string; }",
       sz=11, color=RGBColor(0xA5, 0xB4, 0xFC))

    num(sl, 13)

def s14_auth(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Авторизация: Supabase Auth + JWT",
        subtitle="Токен валидируется на бэкенде через supabase.auth.get_user() — без ручного decode")

    flow = [
        (INDIGO,  "Регистрация /\nВход",       "POST /api/auth/register\nили /login"),
        (TEAL,    "Supabase Auth\nвыдаёт JWT", "access_token подписан\nsupabase_jwt_secret"),
        (VIOLET,  "Фронтенд\nсохраняет токен", "localStorage\n.setItem('access_token')"),
        (AMBER,   "Axios interceptor\nдобавляет заголовок", "Authorization:\nBearer <token>"),
        (GREEN,   "FastAPI\nvalidates",        "get_current_user()\n→ supabase.auth.get_user()"),
    ]

    cw = Inches(2.32)
    ch = Inches(2.5)
    gx = Inches(0.22)
    y0 = Inches(1.5)

    for i, (col, title, detail) in enumerate(flow):
        x = Inches(0.4) + i * (cw + gx)
        step_box(sl, x, y0, cw, ch, str(i+1), title, detail, col)
        if i < 4:
            arrow_h(sl, x + cw + Inches(0.0), y0 + Inches(1.1))

    # Код get_current_user
    code_block(sl, Inches(0.5), Inches(4.2), Inches(6.0), Inches(3.05),
        "async def get_current_user(\n"
        "    token: str = Depends(oauth2_scheme),\n"
        "    db = Depends(get_db)\n"
        "):\n"
        "    res = supabase.auth.get_user(token)\n"
        "    if not res.user:\n"
        "        raise HTTPException(401)\n"
        "    return {'user_id': res.user.id,\n"
        "            'email':   res.user.email}")

    # Axios interceptor
    code_block(sl, Inches(6.7), Inches(4.2), Inches(6.1), Inches(3.05),
        "// api.ts — единственное место с токеном\n"
        "const api = axios.create({ baseURL: BASE_URL });\n\n"
        "api.interceptors.request.use((config) => {\n"
        "  const token = localStorage\n"
        "    .getItem('access_token');\n"
        "  if (token)\n"
        "    config.headers.Authorization =\n"
        "      `Bearer ${token}`;\n"
        "  return config;\n"
        "});")

    num(sl, 14)

def s15_db(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "База данных: 8 таблиц, PostgreSQL через Supabase")

    tables = [
        ("user_profiles",     INDIGO,  "user_id (PK), display_name, bio",
         "Профиль учителя"),
        ("checks",            TEAL,    "id, user_id, pupil_id, function_id,\nfolder_id, errors (JSON), score",
         "Результаты проверок"),
        ("functions",         VIOLET,  "id, user_id, name, system_prompt,\nuser_template, is_default, is_published",
         "Функции проверки"),
        ("function_versions", AMBER,   "id, function_id, version_number,\nsystem_prompt, change_note",
         "История версий"),
        ("pupils",            GREEN,   "id, user_id, name\nUNIQUE(user_id, name)",
         "Ученики"),
        ("groups",            RED,     "id, user_id, name, description",
         "Группы/классы"),
        ("pupil_groups",      GRAY,    "pupil_id (FK), group_id (FK)",
         "Ученик ↔ Группа M:M"),
        ("folders",           INDIGO,  "id, user_id, name, description",
         "Папки для работ"),
    ]

    cw = Inches(5.9)
    ch = Inches(0.95)
    gy = Inches(0.12)

    for i, (name, col, fields, desc) in enumerate(tables):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.5) + col_i * (cw + Inches(0.4))
        y = Inches(1.42) + row_i * (ch + gy)
        box(sl, x, y, cw, ch, WHITE)
        box(sl, x, y, cw, Pt(5), col)
        box(sl, x, y, Inches(0.06), ch, col)
        tb(sl, x + Inches(0.16), y + Inches(0.06),
           Inches(1.5), Inches(0.35), name, sz=12, bold=True, color=col)
        tb(sl, x + Inches(0.16), y + Inches(0.44),
           cw - Inches(0.22), Inches(0.42),
           fields, sz=10, color=GRAY)
        tag(sl, x + Inches(1.75), y + Inches(0.06),
            Inches(4.0), Inches(0.32), desc, col)

    # JSON поле
    box(sl, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.82), INDIGO_LL)
    tb(sl, Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.65),
       "💡  checks.errors — JSON-колонка: хранит Python-список объектов ошибок напрямую. "
       "SQLAlchemy автоматически\nсериализует/десериализует. "
       "Нет JOIN — все данные проверки читаются одним запросом.",
       sz=12, color=INDIGO)

    num(sl, 15)

def s16_tech_stack(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Технологический стек: обоснование выбора")

    rows = [
        (INDIGO,  "FastAPI 0.115",    "Бэкенд",       "Нативный async + SSE, автодокументация Swagger, DI через Depends()"),
        (INDIGO,  "Python 3.11+",     "Язык",         "Единственный язык с нативной поддержкой PaddleOCR, Tesseract и LLM"),
        (TEAL,    "SQLAlchemy 2.0",   "ORM",          "Полный async API, изоляция SQL от бизнес-логики, декларативные модели"),
        (TEAL,    "Supabase",         "БД + Auth",    "Управляемый PostgreSQL + готовая JWT-авторизация, бесплатный хостинг"),
        (VIOLET,  "PaddleOCR 2.8.1", "OCR",          "Наивысшая точность для рукописи, модель SVTR + детектор DB"),
        (VIOLET,  "Tesseract 5.x",   "OCR резерв",   "Компенсирует слабость PaddleOCR на печатном тексте"),
        (AMBER,   "MarkItDown",       "Документы",    "Единый API для PDF/DOCX/PPTX/TXT от Microsoft, сохраняет структуру"),
        (AMBER,   "Ollama qwen2.5:7b","LLM",         "Локальный запуск — данные не утекают; 7B оптимален на CPU"),
        (GREEN,   "React 19",         "Фронтенд",     "Стандарт отрасли, компонентный подход, наибольшая экосистема"),
        (GREEN,   "TypeScript ~6",    "Типизация",    "Обнаруживает ошибки типов на этапе сборки, улучшает IDE"),
        (RED,     "Vite 8",           "Сборщик",      "Мгновенный старт через ES-модули, Rolldown на Rust — быстрее Webpack"),
        (RED,     "Tailwind CSS v4",  "Стили",        "Utility-first: стили прямо в JSX, v4 через Vite-плагин без конфига"),
    ]

    cx = [Inches(0.5), Inches(2.35), Inches(3.9), Inches(5.65)]
    cw_ = [Inches(1.78), Inches(1.48), Inches(1.68), Inches(7.25)]
    hdr_t = ["Технология", "Версия", "Назначение", "Почему выбрано"]

    y_h = Inches(1.4)
    for cw2, cx2, ht in zip(cw_, cx, hdr_t):
        box(sl, cx2, y_h, cw2 - Inches(0.04), Inches(0.48), INDIGO)
        tb(sl, cx2 + Inches(0.08), y_h + Inches(0.07),
           cw2 - Inches(0.14), Inches(0.36),
           ht, sz=12, bold=True, color=WHITE)

    for r, (col, tech, ver, reason) in enumerate(rows):
        y = Inches(1.93) + r * Inches(0.43)
        bg_c = WHITE if r % 2 == 0 else GRAY_L
        cells = [tech, ver, reason[:20], reason]
        for j, (cw2, cx2, cell) in enumerate(zip(cw_, cx, [tech, tech, tech, reason])):
            box(sl, cx2, y, cw2 - Inches(0.04), Inches(0.40), bg_c)
            if j == 0:
                box(sl, cx2, y, Pt(4), Inches(0.40), col)
                tb(sl, cx2 + Inches(0.1), y + Inches(0.06),
                   cw2 - Inches(0.16), Inches(0.32),
                   tech, sz=11, bold=True, color=col)
            elif j == 3:
                tb(sl, cx2 + Inches(0.08), y + Inches(0.05),
                   cw2 - Inches(0.14), Inches(0.32),
                   reason, sz=10, color=DARK)
            else:
                val = ver if j == 1 else cells[2]
                if j == 2:
                    tb(sl, cx2 + Inches(0.08), y + Inches(0.05),
                       cw2 - Inches(0.14), Inches(0.32),
                       ver, sz=10, color=GRAY)

    num(sl, 16)

def s17_patterns(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Паттерны проектирования",
        subtitle="Strategy + Repository + Dependency Injection")

    pats = [
        (INDIGO, "Стратегия (Strategy)",
         "ЗАЧЕМ:\nЛегко добавить новый OCR-движок\nбез изменения существующего кода\n\n"
         "КАК:\nOCRStrategy — абстрактный класс\nPaddleStrategy, TesseractStrategy — реализации\n"
         "HybridOCR — агрегирует обе (◆)",
         "class OCRStrategy(ABC):\n"
         "  @abstractmethod\n"
         "  def recognize(self,\n"
         "    image_path: str) -> OcrResult: ...\n"
         "  @abstractmethod\n"
         "  def get_name(self) -> str: ..."),
        (TEAL, "Репозиторий (Repository)",
         "ЗАЧЕМ:\nСервисы не знают про SQL.\nЛегко сменить СУБД.\nКод тестируем.\n\n"
         "КАК:\nCheckRepository, GroupRepository,\nFunctionRepository и другие.\n"
         "Единственное место с select/insert/update",
         "class CheckRepository:\n"
         "  async def create(self, data) -> Check:\n"
         "    obj = Check(**data)\n"
         "    self.db.add(obj)\n"
         "    await self.db.commit()\n"
         "    return obj"),
        (VIOLET, "Инъекция зависимостей (DI)",
         "ЗАЧЕМ:\nКомпоненты не создают\nзависимости самостоятельно.\nУпрощает тестирование.\n\n"
         "КАК:\nFastAPI Depends() создаёт сессию БД\nи текущего пользователя\nдля каждого запроса",
         "@router.get('/history')\nasync def history(\n"
         "  db = Depends(get_db),\n"
         "  user = Depends(get_current_user)\n"
         "):\n  return await CheckRepository(db)\n"
         "           .get_by_user(user['user_id'])"),
    ]

    cw = Inches(4.0)
    ch = Inches(5.55)
    gx = Inches(0.2)
    y0 = Inches(1.42)

    for i, (col, title, desc, code) in enumerate(pats):
        x = Inches(0.38) + i * (cw + gx)
        box(sl, x, y0, cw, ch, WHITE)
        box(sl, x, y0, cw, Pt(5), col)
        box(sl, x, y0, cw, Inches(0.55), col)
        tb(sl, x + Inches(0.12), y0 + Inches(0.08),
           cw - Inches(0.2), Inches(0.42),
           title, sz=14, bold=True, color=WHITE)
        tb(sl, x + Inches(0.12), y0 + Inches(0.68),
           cw - Inches(0.2), Inches(2.22),
           desc, sz=11, color=GRAY)
        box(sl, x + Inches(0.08), y0 + Inches(2.98),
            cw - Inches(0.16), Inches(2.42), NEAR_BLACK)
        tb(sl, x + Inches(0.18), y0 + Inches(3.06),
           cw - Inches(0.28), Inches(2.28),
           code, sz=10, color=RGBColor(0xA5, 0xB4, 0xFC))

    num(sl, 17)

def s18_usecase_classdiag(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "UML-диаграммы проекта")

    # Use Case
    box(sl, Inches(0.4), Inches(1.42), Inches(6.1), Inches(5.8), GRAY_L)
    tb(sl, Inches(0.55), Inches(1.52), Inches(5.8), Inches(0.38),
       "Диаграмма вариантов использования", sz=13, bold=True, color=INDIGO)
    box(sl, Inches(0.4), Inches(1.42), Pt(5), Inches(5.8), INDIGO)
    tb(sl, Inches(2.0), Inches(3.6), Inches(3.0), Inches(0.5),
       "[Вставить изображение\nuse_case_diagram]",
       sz=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    stats = [("30", "вариантов\nиспользования"), ("1", "актор\n(Учитель)"), ("5", "разделов")]
    for i, (v, l) in enumerate(stats):
        y = Inches(5.25) + i * 0
        x = Inches(0.55) + i * Inches(1.92)
        box(sl, x, Inches(5.2), Inches(1.75), Inches(0.88), INDIGO_L)
        tb(sl, x + Inches(0.08), Inches(5.26), Inches(0.6), Inches(0.75),
           v, sz=28, bold=True, color=INDIGO, align=PP_ALIGN.CENTER)
        tb(sl, x + Inches(0.72), Inches(5.3), Inches(0.95), Inches(0.72),
           l, sz=11, color=DARK)

    # Class diagram
    box(sl, Inches(6.75), Inches(1.42), Inches(6.15), Inches(5.8), GRAY_L)
    tb(sl, Inches(6.9), Inches(1.52), Inches(5.8), Inches(0.38),
       "Диаграмма классов", sz=13, bold=True, color=VIOLET)
    box(sl, Inches(6.75), Inches(1.42), Pt(5), Inches(5.8), VIOLET)
    tb(sl, Inches(8.35), Inches(3.6), Inches(3.0), Inches(0.5),
       "[Вставить изображение\nclass_diagram]",
       sz=13, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    key_pts = [
        "OCRStrategy — абстрактный класс (контракт)",
        "PaddleStrategy, TesseractStrategy — реализации",
        "HybridOCR ◆→ OCRStrategy (композиция)",
        "CheckService → LLMService, FunctionRepo",
    ]
    for i, pt in enumerate(key_pts):
        tb(sl, Inches(6.9), Inches(5.22) + i * Inches(0.35),
           Inches(5.85), Inches(0.32),
           "▸  " + pt, sz=11, color=DARK)

    num(sl, 18)

def s19_monorepo(prs):
    sl = blank_slide(prs)
    fill_bg(sl, SLIDE_BG)
    hdr(sl, "Структура проекта: монорепозиторий",
        subtitle="Две независимые части в одном git-репозитории, развёртывание через Docker")

    # Дерево папок
    tree = [
        ("russian-handwriting-checker-new/", 0, INDIGO),
        ("├── frontend/",                   1, TEAL),
        ("│   ├── src/",                    2, TEAL),
        ("│   │   ├── pages/",              3, GRAY),
        ("│   │   ├── components/",         3, GRAY),
        ("│   │   ├── context/",            3, GRAY),
        ("│   │   ├── api.ts",              3, GRAY),
        ("│   │   └── types.ts",            3, GRAY),
        ("│   └── package.json",            2, GRAY),
        ("├── python-service/",             1, VIOLET),
        ("│   └── src/",                    2, VIOLET),
        ("│       ├── api/routes/",         3, GRAY),
        ("│       ├── services/",           3, GRAY),
        ("│       ├── repositories/",       3, GRAY),
        ("│       ├── models/",             3, GRAY),
        ("│       ├── strategies/",         3, GRAY),
        ("│       └── converters/",         3, GRAY),
        ("├── docker-compose.yml",          1, GREEN),
        ("└── .env",                        1, AMBER),
    ]

    box(sl, Inches(0.5), Inches(1.42), Inches(5.2), Inches(5.8), NEAR_BLACK)
    y_t = Inches(1.55)
    for text, indent, col in tree:
        tb(sl, Inches(0.65) + indent * Inches(0.25), y_t,
           Inches(4.9), Inches(0.3),
           text, sz=11, color=col)
        y_t += Inches(0.28)

    # Описание частей
    parts = [
        (TEAL,   "frontend/",          "React + TypeScript + Vite\nТолько UI, нет бизнес-логики\nСтатика, собирается в dist/"),
        (VIOLET, "python-service/",    "FastAPI + SQLAlchemy\nBCR + LLM + авторизация\nВся бизнес-логика здесь"),
        (GREEN,  "docker-compose.yml", "Запускает оба сервиса\nРаздаёт порты и переменные\nПростой деплой: docker-compose up"),
        (AMBER,  ".env",               "DATABASE_URL, Supabase keys\nOllama URL и модель\nБезопасно: не входит в git"),
    ]

    cw = Inches(3.55)
    ch = Inches(1.38)
    gy = Inches(0.14)

    for i, (col, title, desc) in enumerate(parts):
        x = Inches(5.9)
        y = Inches(1.42) + i * (ch + gy)
        box(sl, x, y, cw, ch, WHITE)
        box(sl, x, y, cw, Pt(5), col)
        box(sl, x, y, Inches(0.06), ch, col)
        tb(sl, x + Inches(0.16), y + Inches(0.08),
           cw - Inches(0.22), Inches(0.38),
           title, sz=13, bold=True, color=col)
        tb(sl, x + Inches(0.16), y + Inches(0.5),
           cw - Inches(0.22), Inches(0.78),
           desc, sz=11, color=GRAY)

    # Docker команда
    code_block(sl, Inches(5.9), Inches(7.0 - Inches(1.35).inches), Inches(7.0), Inches(0.85),
        "docker-compose up --build    # запуск обоих сервисов\n"
        "# frontend → :5173   |   python-service → :8000")

    num(sl, 19)

def s20_results(prs):
    sl = blank_slide(prs)
    fill_bg(sl, INDIGO)
    box(sl, 0, Inches(5.1), W, Inches(2.4), INDIGO_D)
    box(sl, 0, 0, Inches(0.1), H, WHITE)

    tb(sl, Inches(0.6), Inches(0.28), Inches(12.0), Inches(0.75),
       "Результаты работы", sz=32, bold=True, color=WHITE)
    divider(sl, Inches(0.6), Inches(1.1), Inches(10), RGBColor(0xC7, 0xD2, 0xFE))

    results = [
        ("🖊", "Гибридный OCR",       "PaddleOCR + Tesseract, автовыбор лучшего результата"),
        ("📄", "Документы",           "MarkItDown: PDF, DOCX, PPTX, TXT с сохранением структуры"),
        ("🤖", "5 типов проверки",    "Орфография, ОГЭ 9.3, ЕГЭ, итоговое сочинение, свободная"),
        ("⚡", "Стриминг SSE",        "Ответ ИИ появляется в реальном времени, не надо ждать"),
        ("🔐", "JWT-авторизация",     "Supabase Auth, изоляция данных по user_id"),
        ("👨‍🎓", "Управление данными", "Ученики, группы, папки, история с фильтрами"),
        ("⚙️", "Галерея функций",     "Публикация, копирование, версионирование промптов"),
        ("🐳", "Docker",              "Готово к деплою на любом сервере"),
    ]

    for i, (icon, title, desc) in enumerate(results):
        col_i = i % 2
        row_i = i // 2
        x = Inches(0.55) + col_i * Inches(6.35)
        y = Inches(1.28) + row_i * Inches(0.92)
        box(sl, x, y, Inches(6.1), Inches(0.78), RGBColor(0x3B, 0x38, 0xA0))
        tb(sl, x + Inches(0.1), y + Inches(0.1),
           Inches(0.45), Inches(0.58), icon, sz=22)
        tb(sl, x + Inches(0.6), y + Inches(0.08),
           Inches(1.5), Inches(0.35),
           title, sz=13, bold=True, color=WHITE)
        tb(sl, x + Inches(0.6), y + Inches(0.42),
           Inches(5.35), Inches(0.3),
           desc, sz=11, color=INDIGO_L)

    tb(sl, Inches(0.6), Inches(5.28), Inches(12.0), Inches(0.62),
       "Все поставленные задачи выполнены в полном объёме.",
       sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb(sl, Inches(0.6), Inches(5.98), Inches(12.0), Inches(0.62),
       "Спасибо за внимание!  Готов(а) ответить на вопросы.",
       sz=22, bold=True, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)

    num(sl, 20)


# ── Главная функция ────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    s01_title(prs)
    s02_problem(prs)
    s03_goals(prs)
    s04_competitors(prs)
    s05_features(prs)
    s06_check_types(prs)
    s07_architecture(prs)
    s08_ocr_pipeline(prs)
    s09_streaming(prs)
    s10_upload_algo(prs)
    s11_check_service(prs)
    s12_functions(prs)
    s13_frontend(prs)
    s14_auth(prs)
    s15_db(prs)
    s16_tech_stack(prs)
    s17_patterns(prs)
    s18_usecase_classdiag(prs)
    s19_monorepo(prs)
    s20_results(prs)

    out = "RusYazyk_AI_v2.pptx"
    prs.save(out)
    print(f"Сохранено: {out}  ({len(prs.slides)} слайдов)")


if __name__ == "__main__":
    main()
